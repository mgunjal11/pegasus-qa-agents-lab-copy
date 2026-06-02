#!/usr/bin/env python3
"""Generate MSC Dev Code and QA Test Coverage Validator deck — WBD QBR brand + executive presentation visuals."""

from __future__ import annotations

import re
import sys
from pathlib import Path

SCRIPTS_DIR = Path(__file__).resolve().parent
if str(SCRIPTS_DIR) not in sys.path:
    sys.path.insert(0, str(SCRIPTS_DIR))

from coverage_report_helpers import (
    HEADER_H1_INFO,
    OWNERSHIP_LABEL_INFO,
    QUICK_LINKS_INFO,
    READINESS_ITEM_INFO,
    READINESS_PANEL_INFO,
    SECTION_HEADER_INFO,
    SECTION_LEAD_INFO,
    SUMMARY_GROUP_INFO,
    SUMMARY_METRIC_INFO,
    VERDICT_INFO,
)
from ppt_report_from_html import ReportSection, parse_msc_report_html
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
KEEP_PREFIX_SLIDES = 10
DEFAULT_BASE_PPT = ROOT / "reports" / "MSC-Dev-Code-and-QA-Test-Coverage-Validator-Guide.pptx"
DEFAULT_REPORT_HTML = ROOT / "reports" / "MSC-205625-06-02-2026-17-29-55-IST.html"

# Brand
CORAL = RGBColor(0xFF, 0x5E, 0x4F)
NAVY = RGBColor(0x04, 0x06, 0x6C)
NAVY_LIGHT = RGBColor(0x1E, 0x3A, 0x8A)
GOLD = RGBColor(0xFF, 0xC0, 0x00)
GOLD_DARK = RGBColor(0xD4, 0x9A, 0x00)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
SOFT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
SOFT_GOLD = RGBColor(0xFF, 0xF8, 0xE1)
SOFT_CORAL = RGBColor(0xFF, 0xED, 0xEA)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
BODY = RGBColor(0x33, 0x41, 0x55)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
GREEN_BG = RGBColor(0xDC, 0xFC, 0xE7)
AMBER = RGBColor(0xD9, 0x77, 0x06)
AMBER_BG = RGBColor(0xFF, 0xED, 0xD5)
RED = RGBColor(0xDC, 0x26, 0x26)
RED_BG = RGBColor(0xFE, 0xE2, 0xE2)
TEAL = RGBColor(0x00, 0x96, 0x88)
PILOT_ORANGE = RGBColor(0xE8, 0x78, 0x08)
IMPL_GOLD = RGBColor(0xBC, 0xB8, 0x1E)
PILL_GRAY = RGBColor(0x74, 0x74, 0x74)

FONT = "Calibri"
W = Inches(13.333)
H = Inches(7.5)

# Latest report snapshots — refreshed from reports/{KEY}-*-IST.html at build time
REPORT_DEVELOPER = "Mayur Gunjal"
# Featured validation run for example slides (report UI, metrics, traceability, matrix highlight).
PRIMARY_EXAMPLE_KEY = "MSC-205625"

DECK_SLIDE_TIPS: dict[str, str] = {
    "Executive summary": "One automated readiness view before MSC release sign-off.",
    "Who uses it — four personas": "Same slash command; each role reads a different slice of the HTML report.",
    "Automated workflow": "End-to-end --auto --write: one MCP Jira batch + batched Python scripts.",
    "Solution architecture": "Cursor IDE, Atlassian MCP, gh/prefetch scripts, and HTML builder.",
    "HTML report — 8 sections": "Mirrors live report structure.",
    "Report UX — info icons & tooltips": "HTML: apply_report_ui_enhancements() layout v8.",
    "§1 Coverage summary — 8 metric cards": SECTION_HEADER_INFO.get("Coverage summary", ""),
    "§3 Test plan validation — Jira Excel + LADR alignment": SECTION_HEADER_INFO.get(
        "Attached test plan validation", ""
    ),
    "§4–§5 Dev vs QA ownership & traceability": (
        f"{SECTION_HEADER_INFO.get('Dev vs QA test ownership', '')} "
        f"{SECTION_LEAD_INFO}"
    ).strip(),
    "Confluence LADR requirements from Jira": (
        "Fetches LADR/design requirements when linked; quick links exclude grooming/deployment pages."
    ),
    "Latest report matrix — 8-card coverage summary": (
        "Snapshot from newest HTML reports in reports/ — refreshed when the deck is built."
    ),
    "Proven MSC outcomes": "Recent MSC-204417, MSC-205625, MSC-195138, MSC-212571 validation runs.",
    "The challenge": "Why siloed Jira, GitHub, Excel, and Confluence block confident releases.",
    "Idea generation": "Validator unifies evidence into one Pass / Pass with gaps / Fail verdict.",
}

WORKFLOW_STEP_TIPS: dict[str, str] = {
    "Jira": "Parallel MCP: getJiraIssue + remote wiki links + attachments.",
    "Confluence": "fetch_confluence_requirements.py — LADR ESS when linked from Jira.",
    "Test plan": "fetch_jira_testplan.py — Excel attachment, GWT, Mascot/SIT evidence columns.",
    "GitHub": "prefetch_coverage_inputs.py — all PR URLs in one shell invocation.",
    "Map": "map_requirements_to_diff.py — R1…Rn vs diff; qaScope: none when dev-covered.",
    "Report": "build_coverage_report.py + apply_report_ui_enhancements() tooltips v8.",
}

PERSONA_TIPS: dict[str, str] = {
    "Dev Lead": "Dev code %, dev test %, PR diff evidence, CI line coverage.",
    "QA Lead": "Test plan AC %, QA scope remaining, §4 execute list, open gaps.",
    "Release Mgr": "Verdict, release readiness score, open gaps by severity.",
    "Test Designer": "GWT quality, LADR ↔ test case traceability, plan gaps.",
}
REPORT_MATRIX = {
    "MSC-212571": {
        "type": "Story · In QA",
        "generated": "06-02-2026 17:33 IST",
        "verdict": "Pass with gaps",
        "dev_code_pct": "100.0%",
        "dev_tests_pct": "5.6%",
        "req_mapped": "6/9 AC in test plan",
        "testplan_ac_pct": "66.7%",
        "qa_remaining": "9 item(s)",
        "open_gaps": "0 High · 11 Med",
        "ci_line_pct": "99.2%",
        "ci_branch_pct": "99.2%",
        "pr_note": "PR #99 partner-config OPEN · #148 transmogrifier MERGED",
        "testplan_note": "DirecTV Scenarios · 28/28 GWT",
        "ladr_note": "—",
        "report_file": "MSC-212571-06-02-2026-17-33-50-IST.html",
        "summary_short": "QA Testing DirecTV features",
    },
    "MSC-205625": {
        "type": "Bug · Ready for Release",
        "generated": "06-02-2026 17:29 IST",
        "verdict": "Pass with gaps",
        "dev_code_pct": "87.5%",
        "dev_tests_pct": "83.3%",
        "req_mapped": "7/9 AC in test plan",
        "testplan_ac_pct": "77.8%",
        "qa_remaining": "2 item(s)",
        "open_gaps": "0 High · 1 Med",
        "ci_line_pct": "95.3%",
        "ci_branch_pct": "94.5%",
        "pr_note": "PR #161 pick-genie · #195 encode-monitor MERGED",
        "testplan_note": "Domino Inc as full · 5/5 GWT",
        "ladr_note": "LADR quick link only (no grooming wiki)",
        "report_file": "MSC-205625-06-02-2026-17-29-55-IST.html",
        "summary_short": "PFT Clear passport incremental-as-full (SIT)",
    },
    "MSC-204417": {
        "type": "Story · Ready for Release",
        "generated": "06-02-2026 17:26 IST",
        "verdict": "Pass with gaps",
        "dev_code_pct": "100.0%",
        "dev_tests_pct": "33.3%",
        "req_mapped": "15/27 AC in test plan",
        "testplan_ac_pct": "55.6%",
        "qa_remaining": "2 item(s)",
        "open_gaps": "0 High · 2 Med",
        "ci_line_pct": "NA",
        "ci_branch_pct": "NA",
        "pr_note": "No PR · develop branch (pegasus-ess)",
        "testplan_note": "Caption Monitoring · 12/12 GWT",
        "ladr_note": "LADR Captions Delivery Visibility only",
        "report_file": "MSC-204417-06-02-2026-17-26-52-IST.html",
        "summary_short": "V2 caption messaging for Monitor",
    },
    "MSC-195138": {
        "type": "Story · Done",
        "generated": "06-02-2026 17:32 IST",
        "verdict": "Pass with gaps",
        "dev_code_pct": "83.3%",
        "dev_tests_pct": "0.0%",
        "req_mapped": "2/3 AC in test plan",
        "testplan_ac_pct": "66.7%",
        "qa_remaining": "3 item(s)",
        "open_gaps": "0 High · 2 Med",
        "ci_line_pct": "77.7%",
        "ci_branch_pct": "77.7%",
        "pr_note": "PR #22 reps · #75 texttransform MERGED",
        "testplan_note": "FF Race Scenarios · 11/11 GWT",
        "ladr_note": "—",
        "report_file": "MSC-195138-06-02-2026-17-32-47-IST.html",
        "summary_short": "FF2.0 messaging race conditions",
    },
}


def _metric_after_label(html: str, label: str) -> str:
    pattern = re.escape(label) + r"</div>.*?metric-value\">([^<]+)"
    m = re.search(pattern, html, re.S | re.I)
    return (m.group(1).strip() if m else "") or "—"


def _notes_from_report_html(html: str) -> dict[str, str]:
    """Derive pr_note, testplan_note, ladr_note from report header blocks."""
    out: dict[str, str] = {}
    ql = re.search(r'class="quick-links[^"]*">(.*?)</div>', html, re.S)
    block = ql.group(1) if ql else ""
    prs = re.findall(r"github\.com/[^/]+/([^/]+)/pull/(\d+)", block)
    if prs:
        parts = [f"#{n} {repo}" for repo, n in prs[:2]]
        if len(prs) > 2:
            parts.append(f"+{len(prs) - 2} PR")
        out["pr_note"] = "PR " + " · ".join(parts)
    elif re.search(r"branch compare|develop branch|No linked PR", html, re.I):
        repo_m = re.search(r"pegasus-[\w-]+", html)
        repo = repo_m.group(0) if repo_m else "repo"
        out["pr_note"] = f"No PR · develop branch ({repo})"
    else:
        out["pr_note"] = "No linked PR"
    tp = re.search(r"<strong>Test plan</strong>.*?—\s*([^<]{20,200})", html, re.S)
    if tp:
        snippet = re.sub(r"\s+", " ", tp.group(1)).strip()
        sheet_m = re.search(r"sheet\s+(\S+)|tab\s+(\S+)", snippet, re.I)
        scen_m = re.search(r"(\d+)\s+scenarios?", snippet, re.I)
        gwt_m = re.search(r"(\d+)/(\d+)\s+full\s+Given", html, re.I)
        parts = []
        if sheet_m:
            parts.append((sheet_m.group(1) or sheet_m.group(2) or "").strip())
        if scen_m:
            parts.append(f"{scen_m.group(1)} scenarios")
        if gwt_m:
            parts.append(f"{gwt_m.group(1)}/{gwt_m.group(2)} GWT")
        out["testplan_note"] = " · ".join(parts) if parts else snippet[:60]
    ladr_links = re.findall(r"wiki[^\"]*LADR[^\"]*|LADR[^\"]*wiki", block, re.I)
    grooming_in_ql = bool(re.search(r"grooming|deployment|go live|pvc go", block, re.I))
    if ladr_links and not grooming_in_ql:
        out["ladr_note"] = "LADR Confluence in quick links only"
    elif grooming_in_ql:
        out["ladr_note"] = "—"
    else:
        conf = re.search(r"Confluence / LADR</strong>.*?—\s*([^<]+)", html, re.S)
        if conf and "loaded" in conf.group(1).lower():
            out["ladr_note"] = "LADR requirements in cache"
        else:
            out["ladr_note"] = "—"
    return out


def refresh_report_matrix_from_html(reports_dir: Path | None = None) -> str | None:
    """Load newest HTML report per MSC key; return key used as LATEST_EXAMPLE."""
    base = reports_dir or (ROOT / "reports")
    if not base.is_dir():
        return None
    by_key: dict[str, Path] = {}
    for path in base.glob("MSC-*.html"):
        if "Guide" in path.name:
            continue
        m = re.match(r"(MSC-\d+)-", path.name)
        if not m:
            continue
        key = m.group(1)
        prev = by_key.get(key)
        if not prev or path.stat().st_mtime > prev.stat().st_mtime:
            by_key[key] = path

    latest_key: str | None = None
    latest_mtime = 0.0
    for key, path in by_key.items():
        html = path.read_text(encoding="utf-8", errors="replace")
        entry = REPORT_MATRIX.setdefault(
            key,
            {
                "type": "Story",
                "generated": "",
                "verdict": "Pass with gaps",
                "dev_code_pct": "—",
                "dev_tests_pct": "—",
                "req_mapped": "—",
                "testplan_ac_pct": "—",
                "qa_remaining": "—",
                "open_gaps": "—",
                "ci_line_pct": "NA",
                "ci_branch_pct": "NA",
                "pr_note": "—",
                "testplan_note": "—",
                "ladr_note": "—",
                "report_file": path.name,
                "summary_short": key,
            },
        )
        entry["report_file"] = path.name
        gen = re.search(r"Generated:.*?(\d{4})-(\d{2})-(\d{2}) (\d{2}:\d{2})", html)
        if gen:
            entry["generated"] = f"{gen.group(2)}-{gen.group(3)}-{gen.group(1)} {gen.group(4)} IST"
        status_m = re.search(r"Status:.*?>([^<]+)</span>\s*&nbsp;\|", html)
        type_m = re.search(r"Type:.*?>([^<]+)</span>\s*&nbsp;\|", html)
        if status_m and type_m:
            entry["type"] = f"{type_m.group(1).strip()} · {status_m.group(1).strip()}"
        title_m = re.search(r"Coverage validation: MSC-\d+ — ([^<]+)<", html)
        if title_m:
            entry["summary_short"] = title_m.group(1).strip()[:80]
        if "verdict-fail" in html:
            entry["verdict"] = "Fail"
        elif "verdict-pass-gaps" in html:
            entry["verdict"] = "Pass with gaps"
        elif "verdict-pass" in html:
            entry["verdict"] = "Pass"
        entry["dev_code_pct"] = _metric_after_label(html, "Dev code coverage") or entry["dev_code_pct"]
        entry["dev_tests_pct"] = _metric_after_label(html, "Dev unit / integration test coverage") or entry["dev_tests_pct"]
        entry["req_mapped"] = _metric_after_label(html, "Requirements mapped") or entry["req_mapped"]
        entry["testplan_ac_pct"] = _metric_after_label(html, "Test plan acceptance criteria coverage") or entry["testplan_ac_pct"]
        entry["qa_remaining"] = _metric_after_label(html, "QA scope remaining") or entry["qa_remaining"]
        entry["open_gaps"] = _metric_after_label(html, "Open gaps") or entry["open_gaps"]
        entry["ci_line_pct"] = _metric_after_label(html, "CI line coverage") or entry["ci_line_pct"]
        entry["ci_branch_pct"] = _metric_after_label(html, "CI branch coverage") or entry["ci_branch_pct"]
        for k, v in _notes_from_report_html(html).items():
            if v and v != "—":
                entry[k] = v
        mtime = path.stat().st_mtime
        if mtime > latest_mtime:
            latest_mtime = mtime
            latest_key = key
    return latest_key


def _set_latest_example(key: str | None) -> None:
    global LATEST_EXAMPLE
    k = key or PRIMARY_EXAMPLE_KEY
    if k not in REPORT_MATRIX:
        k = next(iter(REPORT_MATRIX))
    LATEST_EXAMPLE = {**REPORT_MATRIX[k], "key": k, "summary_short": REPORT_MATRIX[k].get("summary_short", k)}


_set_latest_example(PRIMARY_EXAMPLE_KEY)


def _delete_slide(prs: Presentation, index: int) -> None:
    """Remove slide at *index* (0-based) from an open presentation."""
    sld_id_lst = prs.slides._sldIdLst
    sld_id = sld_id_lst[index]
    r_id = sld_id.rId
    prs.part.drop_rel(r_id)
    sld_id_lst.remove(sld_id)


class Deck:
    def __init__(self, prs: Presentation | None = None, *, footer_start: int = 0):
        if prs is None:
            self.prs = Presentation()
            self.prs.slide_width = W
            self.prs.slide_height = H
        else:
            self.prs = prs
        self._n = footer_start

    def blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def footer(self, slide):
        self._n += 1
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.92), W, Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = LIGHT_GRAY
        line.line.fill.background()
        lb = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(9), Inches(0.35))
        lb.text_frame.text = (
            f"©LTM  |  Privileged and Confidential  |  MSC Dev Code and QA Test Coverage Validator  |  "
            f"Developed by {REPORT_DEVELOPER}"
        )
        p = lb.text_frame.paragraphs[0]
        p.font.size = Pt(8)
        p.font.name = FONT
        p.font.color.rgb = MUTED
        rb = slide.shapes.add_textbox(Inches(12.35), Inches(7.0), Inches(0.6), Inches(0.35))
        rb.text_frame.text = str(self._n)
        rb.text_frame.paragraphs[0].font.size = Pt(8)
        rb.text_frame.paragraphs[0].font.name = FONT
        rb.text_frame.paragraphs[0].font.color.rgb = MUTED
        rb.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    def _accent_bar(self, slide, top=Inches(0)):
        b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, Inches(0.12), H)
        b.fill.solid()
        b.fill.fore_color.rgb = CORAL
        b.line.fill.background()

    def _tip_for(self, label: str, *maps: dict[str, str]) -> str:
        key = label.strip().rstrip(":")
        for m in maps:
            if key in m and m[key]:
                return m[key]
        return ""

    def _slide_title(self, slide, title: str, subtitle: str | None = None, tip: str | None = None):
        self._accent_bar(slide)
        tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(12.0), Inches(0.65))
        tb.text_frame.text = title
        p = tb.text_frame.paragraphs[0]
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.name = FONT
        p.font.color.rgb = NAVY
        if subtitle:
            st = slide.shapes.add_textbox(Inches(0.55), Inches(0.88), Inches(11.8), Inches(0.32))
            st.text_frame.text = subtitle
            sp = st.text_frame.paragraphs[0]
            sp.font.size = Pt(13)
            sp.font.name = FONT
            sp.font.color.rgb = MUTED

    def _rect(
        self,
        slide,
        l,
        t,
        w,
        h,
        fill,
        text=None,
        font_pt=11,
        bold=False,
        color=WHITE,
        radius=None,
    ):
        kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        sh = slide.shapes.add_shape(kind, l, t, w, h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        sh.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        sh.line.width = Pt(0.5)
        if text is not None:
            p0 = sh.text_frame.paragraphs[0]
            p0.text = text
            p0.font.size = Pt(font_pt)
            p0.font.bold = bold
            p0.font.name = FONT
            p0.font.color.rgb = color
            p0.alignment = PP_ALIGN.CENTER
            sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            sh.text_frame.word_wrap = True
        return sh

    def _icon_circle(self, slide, l, t, size, letter, fill=NAVY):
        o = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, size, size)
        o.fill.solid()
        o.fill.fore_color.rgb = fill
        o.line.fill.background()
        o.text_frame.text = letter
        p = o.text_frame.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.name = FONT
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        o.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        return o

    def _arrow(self, slide, x1, y1, x2, y2):
        c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        c.line.color.rgb = GOLD
        c.line.width = Pt(2.5)
        return c

    def _kpi_card(self, slide, l, t, w, h, value, label, sub, fill, value_color=WHITE, tooltip: str | None = None):
        card = self._rect(slide, l, t, w, h, fill, radius=True)
        v = slide.shapes.add_textbox(l + Inches(0.1), t + Inches(0.12), w - Inches(0.2), Inches(0.5))
        v.text_frame.text = value
        v.text_frame.paragraphs[0].font.size = Pt(28)
        v.text_frame.paragraphs[0].font.bold = True
        v.text_frame.paragraphs[0].font.name = FONT
        v.text_frame.paragraphs[0].font.color.rgb = value_color
        v.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        lb = slide.shapes.add_textbox(l + Inches(0.08), t + Inches(0.58), w - Inches(0.16), Inches(0.38))
        lb.text_frame.text = label
        lb.text_frame.paragraphs[0].font.size = Pt(10)
        lb.text_frame.paragraphs[0].font.bold = True
        lb.text_frame.paragraphs[0].font.name = FONT
        lb.text_frame.paragraphs[0].font.color.rgb = value_color
        lb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        if sub:
            sb = slide.shapes.add_textbox(l + Inches(0.08), t + h - Inches(0.38), w - Inches(0.16), Inches(0.3))
            sb.text_frame.text = sub
            sb.text_frame.paragraphs[0].font.size = Pt(9)
            sb.text_frame.paragraphs[0].font.name = FONT
            sb.text_frame.paragraphs[0].font.color.rgb = value_color
            sb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def title_slide(self):
        s = self.blank()
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = NAVY
        bg.line.fill.background()
        # decorative circles
        for cx, cy, r, alpha in [(Inches(10.5), Inches(-0.8), Inches(3.2), 0.85), (Inches(11.8), Inches(5.5), Inches(2.0), 0.9)]:
            c = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, r, r)
            c.fill.solid()
            c.fill.fore_color.rgb = NAVY_LIGHT
            c.fill.transparency = alpha
            c.line.fill.background()
        band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.55), W, Inches(0.95))
        band.fill.solid()
        band.fill.fore_color.rgb = GOLD
        band.line.fill.background()
        t1 = s.shapes.add_textbox(Inches(0.85), Inches(1.6), Inches(11), Inches(0.9))
        t1.text_frame.text = "Business Assurance"
        t1.text_frame.paragraphs[0].font.size = Pt(22)
        t1.text_frame.paragraphs[0].font.name = FONT
        t1.text_frame.paragraphs[0].font.color.rgb = GOLD
        t2 = s.shapes.add_textbox(Inches(0.85), Inches(2.35), Inches(11), Inches(1.4))
        t2.text_frame.text = "MSC Dev Code and QA Test Coverage Validator"
        t2.text_frame.paragraphs[0].font.size = Pt(40)
        t2.text_frame.paragraphs[0].font.bold = True
        t2.text_frame.paragraphs[0].font.name = FONT
        t2.text_frame.paragraphs[0].font.color.rgb = WHITE
        t3 = s.shapes.add_textbox(Inches(0.85), Inches(3.75), Inches(10), Inches(0.9))
        t3.text_frame.text = "AI-driven release readiness · Jira + Confluence LADR (when linked) + GitHub + Excel test plan in one report"
        t3.text_frame.paragraphs[0].font.size = Pt(16)
        t3.text_frame.paragraphs[0].font.name = FONT
        t3.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
        # KPI strip on title
        metrics = [("4", "Systems unified"), ("8", "Coverage metrics"), ("2–5", "Minutes per run"), ("1", "HTML verdict")]
        cw = Inches(2.85)
        title_kpis = {
            "Systems unified": "Jira, Confluence LADR, GitHub PR(s), and Excel test plan in one run.",
            "Coverage metrics": "Eight summary cards in §1 of the HTML report (dev, plan, CI, gaps).",
            "Minutes per run": "Typical --auto pipeline after caches are warm (2–5 minutes).",
            "HTML verdict": VERDICT_INFO,
        }
        for i, (val, lbl) in enumerate(metrics):
            left = Inches(0.85) + cw * i + Inches(0.12) * i
            self._kpi_card(s, left, Inches(4.85), cw, Inches(1.15), val, lbl, None, NAVY_LIGHT)
        foot = s.shapes.add_textbox(Inches(0.85), Inches(6.72), Inches(11), Inches(0.4))
        foot.text_frame.text = (
            f"Pegasus QA Agents Lab · github.com/mgunjal11/pegasus-qa-agents-lab · "
            f"Developed by {REPORT_DEVELOPER}"
        )
        foot.text_frame.paragraphs[0].font.size = Pt(11)
        foot.text_frame.paragraphs[0].font.name = FONT
        foot.text_frame.paragraphs[0].font.color.rgb = DARK

    def agenda_slide(self):
        s = self.blank()
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHITE
        self.footer(s)
        self._accent_bar(s)
        tb = s.shapes.add_textbox(Inches(0.55), Inches(0.5), Inches(4), Inches(0.9))
        tb.text_frame.text = "Agenda"
        tb.text_frame.paragraphs[0].font.size = Pt(44)
        tb.text_frame.paragraphs[0].font.bold = True
        tb.text_frame.paragraphs[0].font.color.rgb = CORAL
        tb.text_frame.paragraphs[0].font.name = FONT
        items = [
            ("01", "Executive summary", "Value at a glance"),
            ("02", "The challenge", "Fragmented release signals"),
            ("03", "The solution", "Cursor subagent + workflow"),
            ("04", "HTML report deep-dive", "8 sections · metrics · traceability"),
            ("05", "Enablement", "Setup · scripts · permissions"),
            ("06", "Proven outcomes", "MSC case studies"),
        ]
        for i, (num, title, sub) in enumerate(items):
            top = Inches(1.45) + Inches(1.05) * i
            card = self._rect(
                s,
                Inches(0.55),
                top,
                Inches(6.2),
                Inches(0.95),
                LIGHT_GRAY if i % 2 else SOFT_BLUE,
                radius=True,
            )
            card.line.fill.background()
            nb = s.shapes.add_textbox(Inches(0.7), top + Inches(0.12), Inches(0.55), Inches(0.5))
            nb.text_frame.text = num
            nb.text_frame.paragraphs[0].font.size = Pt(22)
            nb.text_frame.paragraphs[0].font.bold = True
            nb.text_frame.paragraphs[0].font.color.rgb = CORAL
            nb.text_frame.paragraphs[0].font.name = FONT
            tt = s.shapes.add_textbox(Inches(1.35), top + Inches(0.1), Inches(4.5), Inches(0.35))
            tt.text_frame.text = title
            tt.text_frame.paragraphs[0].font.size = Pt(16)
            tt.text_frame.paragraphs[0].font.bold = True
            tt.text_frame.paragraphs[0].font.color.rgb = NAVY
            tt.text_frame.paragraphs[0].font.name = FONT
            st = s.shapes.add_textbox(Inches(1.35), top + Inches(0.42), Inches(4.8), Inches(0.3))
            st.text_frame.text = sub
            st.text_frame.paragraphs[0].font.size = Pt(11)
            st.text_frame.paragraphs[0].font.color.rgb = MUTED
            st.text_frame.paragraphs[0].font.name = FONT
        # right visual — mini architecture
        panel = self._rect(s, Inches(7.1), Inches(1.2), Inches(5.7), Inches(5.5), NAVY, radius=True)
        panel.line.fill.background()
        cap = s.shapes.add_textbox(Inches(7.35), Inches(1.4), Inches(5.2), Inches(0.45))
        cap.text_frame.text = "What you get in one run"
        cap.text_frame.paragraphs[0].font.size = Pt(14)
        cap.text_frame.paragraphs[0].font.bold = True
        cap.text_frame.paragraphs[0].font.color.rgb = GOLD
        cap.text_frame.paragraphs[0].font.name = FONT
        stack = [
            "Jira story + AC + LADR links",
            "Confluence LADR requirements (when linked)",
            "Linked PR + CI",
            "Excel test plan (Jira attachment)",
            "HTML readiness report",
        ]
        for j, lbl in enumerate(stack):
            y = Inches(1.85) + Inches(0.78) * j
            self._rect(s, Inches(7.5), y, Inches(4.9), Inches(0.58), NAVY_LIGHT, lbl, 11, True, WHITE, True)
            if j < len(stack) - 1:
                arr = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9.6), y + Inches(0.72), Inches(0.35), Inches(0.22))
                arr.fill.solid()
                arr.fill.fore_color.rgb = GOLD
                arr.line.fill.background()

    def executive_summary(self):
        s = self.blank()
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHITE
        self.footer(s)
        self._slide_title(s, "Executive summary", "One automated readiness view before every MSC release sign-off")
        cards = [
            ("Problem", "Jira, GitHub, Excel test plans, and Confluence LADR (when linked) live in silos — release reviews lack one evidence trail.", CORAL, WHITE),
            ("Solution", "Cursor subagent correlates Jira AC + Confluence LADR requirements → code → dev tests → QA plan → verdict.", NAVY, WHITE),
            ("Outcome", "Leadership opens a timestamped HTML report with 8 quantified metrics and clear QA handoff.", TEAL, WHITE),
        ]
        for i, (h, b, fill, tc) in enumerate(cards):
            left = Inches(0.55) + Inches(4.15) * i
            self._rect(s, left, Inches(1.35), Inches(3.95), Inches(0.5), fill, h, 14, True, tc, True)
            box = s.shapes.add_textbox(left + Inches(0.12), Inches(1.95), Inches(3.7), Inches(1.2))
            box.text_frame.text = b
            box.text_frame.paragraphs[0].font.size = Pt(12)
            box.text_frame.paragraphs[0].font.name = FONT
            box.text_frame.paragraphs[0].font.color.rgb = BODY
            box.text_frame.word_wrap = True
        # bottom impact row
        impacts = [
            ("~35–40%", "Faster triage vs manual traceability", SOFT_CORAL, CORAL),
            ("100%", "AC-to-evidence mapping per run", SOFT_BLUE, NAVY),
            ("Pass / Gaps / Fail", "Explicit release verdict", SOFT_GOLD, GOLD_DARK),
            ("Audit-ready", "Shareable HTML artifact", GREEN_BG, GREEN),
        ]
        for i, (val, lbl, bg, vc) in enumerate(impacts):
            left = Inches(0.55) + Inches(3.12) * i
            self._kpi_card(
                s,
                left,
                Inches(3.85),
                Inches(2.95),
                Inches(1.55),
                val,
                lbl,
                "Per validator run",
                bg,
                vc,
            )
        quote = self._rect(s, Inches(0.55), Inches(5.65), Inches(12.2), Inches(0.85), LIGHT_GRAY, radius=True)
        quote.line.fill.background()
        qt = s.shapes.add_textbox(Inches(0.75), Inches(5.78), Inches(11.8), Inches(0.6))
        qt.text_frame.text = (
            '"We no longer ask five people whether the PR, unit tests, Jira test plan, and LADR design doc align — '
            'the validator answers that in one report." — MSC QA lead workflow'
        )
        qt.text_frame.paragraphs[0].font.size = Pt(12)
        qt.text_frame.paragraphs[0].font.italic = True
        qt.text_frame.paragraphs[0].font.color.rgb = NAVY
        qt.text_frame.paragraphs[0].font.name = FONT


    def section_slide(self, num: str, title: str, tagline: str, tip: str | None = None):
        s = self.blank()
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = NAVY
        bg.line.fill.background()
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(1.2), Inches(4.5), Inches(4.5))
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY_LIGHT
        c.fill.transparency = 0.25
        c.line.fill.background()
        num_b = s.shapes.add_textbox(Inches(0.85), Inches(2.0), Inches(2), Inches(1.2))
        num_b.text_frame.text = num
        num_b.text_frame.paragraphs[0].font.size = Pt(72)
        num_b.text_frame.paragraphs[0].font.bold = True
        num_b.text_frame.paragraphs[0].font.color.rgb = GOLD
        num_b.text_frame.paragraphs[0].font.name = FONT
        tit = s.shapes.add_textbox(Inches(0.85), Inches(3.2), Inches(10), Inches(1.2))
        tit.text_frame.text = title
        tit.text_frame.paragraphs[0].font.size = Pt(36)
        tit.text_frame.paragraphs[0].font.bold = True
        tit.text_frame.paragraphs[0].font.color.rgb = WHITE
        tit.text_frame.paragraphs[0].font.name = FONT
        tg = s.shapes.add_textbox(Inches(0.85), Inches(4.35), Inches(9), Inches(0.6))
        tg.text_frame.text = tagline
        tg.text_frame.paragraphs[0].font.size = Pt(16)
        tg.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
        tg.text_frame.paragraphs[0].font.name = FONT
        self.footer(s)

    def challenge_slide(self):
        s = self.blank()
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHITE
        self.footer(s)
        self._slide_title(s, "The challenge — four silos, zero single view", "Manual traceability before release is slow, inconsistent, and error-prone")
        silos = [
            ("J", "Jira", "Acceptance criteria,\nstory status, attachments", SOFT_BLUE, NAVY),
            ("GH", "GitHub", "PR diffs, checks,\nunit/integration tests", SOFT_GOLD, GOLD_DARK),
            ("XL", "Excel test plan", "Jira attachment ·\nGiven/When/Then · Mascot or IDs", SOFT_CORAL, CORAL),
            ("CF", "Confluence LADR", "Design requirements\nwhen linked from Jira", RGBColor(0xCC, 0xFB, 0xF1), TEAL),
        ]
        col_w = Inches(2.95)
        gap = Inches(0.15)
        for i, (icon, name, desc, bg, ic) in enumerate(silos):
            left = Inches(0.55) + (col_w + gap) * i
            self._rect(s, left, Inches(1.45), col_w, Inches(3.2), bg, radius=True)
            self._icon_circle(s, left + Inches(1.0), Inches(1.75), Inches(0.95), icon, ic)
            nm = s.shapes.add_textbox(left, Inches(2.85), col_w, Inches(0.4))
            nm.text_frame.text = name
            nm.text_frame.paragraphs[0].font.size = Pt(18)
            nm.text_frame.paragraphs[0].font.bold = True
            nm.text_frame.paragraphs[0].font.color.rgb = DARK
            nm.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            nm.text_frame.paragraphs[0].font.name = FONT
            ds = s.shapes.add_textbox(left + Inches(0.12), Inches(3.35), col_w - Inches(0.24), Inches(1.0))
            ds.text_frame.text = desc
            ds.text_frame.paragraphs[0].font.size = Pt(11)
            ds.text_frame.paragraphs[0].font.color.rgb = BODY
            ds.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            ds.text_frame.paragraphs[0].font.name = FONT
        # pain callouts
        pains = [
            "Duplicate QA/dev effort",
            "Missed E2E handoffs",
            "No LADR ↔ test plan traceability",
            "Weak AC ↔ evidence mapping",
        ]
        for i, p in enumerate(pains):
            left = Inches(0.65) + Inches(3.12) * i
            chip = self._rect(s, left, Inches(5.0), Inches(2.95), Inches(0.55), RED_BG, p, 10, True, RED, True)
            chip.line.color.rgb = RED
        self._rect(s, Inches(0.55), Inches(5.75), Inches(12.2), Inches(0.75), NAVY, radius=True).line.fill.background()
        sol = s.shapes.add_textbox(Inches(0.75), Inches(5.9), Inches(11.8), Inches(0.5))
        sol.text_frame.text = "→  MSC Dev Code and QA Test Coverage Validator unifies all four silos into one HTML readiness report with numbered recommended actions"
        sol.text_frame.paragraphs[0].font.size = Pt(13)
        sol.text_frame.paragraphs[0].font.bold = True
        sol.text_frame.paragraphs[0].font.color.rgb = WHITE
        sol.text_frame.paragraphs[0].font.name = FONT

    def idea_generation_slide(self):
        s = self.blank()
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHITE
        self.footer(s)
        self._slide_title(s, "Idea Generation — Cursor AI subagent", "Completed · Pegasus QA Agents Lab · Gen AI transformation journey")
        cols = [
            ("Fragmented readiness", NAVY, "Problem Statement", [
                "AC in Jira, code in GitHub, Excel test plans, and Confluence LADR — four silos, no single view",
                "Release reviews depend on tribal knowledge",
                "Manual traceability is slow before sign-off",
            ]),
            ("Unified validator", GOLD, "Solution", [
                "/msc-dev-code-and-qa-test-coverage-validator {KEY}",
                "One batch: Jira + Confluence LADR (if linked) + PR(s) + Excel test plan",
                "Maps Jira AC + LADR L1…Ln → test cases → code → dev tests",
                "HTML report + LADR traceability table + Pass / Pass with gaps / Fail",
            ]),
            ("Leadership evidence", NAVY, "Benefits", [
                "8 quantified metric cards",
                "Dev vs QA ownership made explicit",
                "Mascot + Given/When/Then + SIT Jobs ID evidence",
                "Audit-friendly downloadable artifact",
                f"Developed by {REPORT_DEVELOPER}",
            ]),
        ]
        col_w = Inches(3.95)
        for i, (title, hf, sec, bullets) in enumerate(cols):
            left = Inches(0.55) + (col_w + Inches(0.22)) * i
            self._rect(s, left, Inches(1.35), col_w, Inches(0.48), hf, title, 13, True, WHITE, True)
            body = self._rect(s, left, Inches(1.83), col_w, Inches(4.35), LIGHT_GRAY, radius=True)
            body.line.fill.background()
            tb = s.shapes.add_textbox(left + Inches(0.12), Inches(1.95), col_w - Inches(0.24), Inches(4.0))
            tf = tb.text_frame
            tf.word_wrap = True
            p0 = tf.paragraphs[0]
            p0.text = f"{sec}:"
            p0.font.bold = True
            p0.font.size = Pt(11)
            p0.font.color.rgb = NAVY
            p0.font.name = FONT
            for b in bullets:
                para = tf.add_paragraph()
                para.text = f"• {b}"
                para.font.size = Pt(10)
                para.font.name = FONT
                para.font.color.rgb = BODY
                para.space_after = Pt(4)
            pills = ["Ideation", "Pilot", "Implementation", "Completed"]
            fills = [PILL_GRAY, PILOT_ORANGE, IMPL_GOLD, GREEN]
            pw = (col_w - Inches(0.1)) / 4
            for j, (lab, fill) in enumerate(zip(pills, fills)):
                sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.05) + pw * j, Inches(5.95), pw - Inches(0.05), Inches(0.22))
                sh.fill.solid()
                sh.fill.fore_color.rgb = fill if lab == "Completed" else PILL_GRAY
                sh.line.fill.background()
                sh.text_frame.text = lab
                sh.text_frame.paragraphs[0].font.size = Pt(7)
                sh.text_frame.paragraphs[0].font.bold = True
                sh.text_frame.paragraphs[0].font.color.rgb = WHITE
                sh.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def personas_slide(self):
        s = self.blank()
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHITE
        self.footer(s)
        self._slide_title(s, "Who uses it — four personas", "Same command, role-specific value from one HTML report")
        personas = [
            ("QA Lead", "Pre-release readiness", "In QA / Ready for Release stories", "QA scope %, test plan gaps", SOFT_BLUE, NAVY),
            ("Developer", "Merge confidence", "Dev-owned AC + unit tests", "Code %, dev test evidence", SOFT_GOLD, GOLD_DARK),
            ("Release Mgr", "Go/no-go snapshot", "Verdict + open gaps", "8 cards, recommended actions", SOFT_CORAL, CORAL),
            ("Test Designer", "Plan alignment", "Excel plan vs Jira AC + LADR", "GWT quality · Mascot or SIT Jobs IDs", GREEN_BG, GREEN),
        ]
        for i, (role, when, focus, metric, bg, accent) in enumerate(personas):
            left = Inches(0.55) + Inches(3.12) * i
            self._rect(s, left, Inches(1.4), Inches(2.95), Inches(4.8), bg, radius=True)
            self._icon_circle(s, left + Inches(1.0), Inches(1.65), Inches(0.95), role[0], accent)
            for yoff, txt, bold, sz in [(2.75, role, True, 14), (3.15, when, False, 11), (3.55, focus, False, 10), (4.2, metric, True, 10)]:
                bx = s.shapes.add_textbox(left + Inches(0.15), Inches(yoff), Inches(2.65), Inches(0.55))
                bx.text_frame.text = txt
                bx.text_frame.paragraphs[0].font.size = Pt(sz)
                bx.text_frame.paragraphs[0].font.bold = bold
                bx.text_frame.paragraphs[0].font.color.rgb = DARK if bold else BODY
                bx.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                bx.text_frame.paragraphs[0].font.name = FONT
        cmd = self._rect(s, Inches(2.5), Inches(6.35), Inches(8.3), Inches(0.5), DARK, radius=True)
        cmd.line.fill.background()
        ct = s.shapes.add_textbox(Inches(2.65), Inches(6.42), Inches(8), Inches(0.38))
        ct.text_frame.text = f"/msc-dev-code-and-qa-test-coverage-validator {PRIMARY_EXAMPLE_KEY}   ·   --auto --write"
        ct.text_frame.paragraphs[0].font.size = Pt(14)
        ct.text_frame.paragraphs[0].font.name = "Consolas"
        ct.text_frame.paragraphs[0].font.color.rgb = GOLD
        ct.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def workflow_slide(self):
        s = self.blank()
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHITE
        self.footer(s)
        self._slide_title(s, "Automated workflow", "Slash command --auto --write · ~2–5 min · one MCP turn + batched shells")
        steps = [
            ("1", "Jira", "Parallel MCP:\nissue + remote\nwiki links"),
            ("2", "Confluence", "LADR ESS or\npassport scenarios"),
            ("3", "Test plan", "Excel attach,\nGWT, Evidence"),
            ("4", "GitHub", "prefetch PR(s)\nor branch compare"),
            ("5", "Map", "map_requirements_\nto_diff.py"),
            ("6", "Report", "build_coverage_\nreport.py + tooltips"),
        ]
        sw = Inches(1.95)
        for i, (num, title, sub) in enumerate(steps):
            left = Inches(0.35) + (sw + Inches(0.18)) * i
            self._rect(
                s,
                left,
                Inches(1.55),
                sw,
                Inches(2.55),
                NAVY if i % 2 == 0 else NAVY_LIGHT,
                radius=True,
            )
            nb = s.shapes.add_textbox(left + Inches(0.15), Inches(1.7), Inches(0.5), Inches(0.45))
            nb.text_frame.text = num
            nb.text_frame.paragraphs[0].font.size = Pt(22)
            nb.text_frame.paragraphs[0].font.bold = True
            nb.text_frame.paragraphs[0].font.color.rgb = GOLD
            nb.text_frame.paragraphs[0].font.name = FONT
            tt = s.shapes.add_textbox(left, Inches(2.25), sw, Inches(0.4))
            tt.text_frame.text = title
            tt.text_frame.paragraphs[0].font.size = Pt(14)
            tt.text_frame.paragraphs[0].font.bold = True
            tt.text_frame.paragraphs[0].font.color.rgb = WHITE
            tt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            tt.text_frame.paragraphs[0].font.name = FONT
            sb = s.shapes.add_textbox(left + Inches(0.1), Inches(2.75), sw - Inches(0.2), Inches(0.9))
            sb.text_frame.text = sub
            sb.text_frame.paragraphs[0].font.size = Pt(10)
            sb.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
            sb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            sb.text_frame.paragraphs[0].font.name = FONT
            if i < len(steps) - 1:
                ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + sw + Inches(0.04), Inches(2.45), Inches(0.22), Inches(0.35))
                ar.fill.solid()
                ar.fill.fore_color.rgb = GOLD
                ar.line.fill.background()
        # cache callout
        self._rect(s, Inches(0.55), Inches(4.25), Inches(12.2), Inches(1.15), SOFT_GOLD, radius=True)
        cb = s.shapes.add_textbox(Inches(0.75), Inches(4.4), Inches(11.8), Inches(0.9))
        cb.text_frame.text = (
            "Cache: reports/.cache/{KEY}-jira|confluence|testplan|prefetch|mapping.json — reuse --from-cache. "
            "Quick links: collect_ladr_page_links() — LADR/design Confluence only (no grooming/deployment remote links). "
            "No linked PR? fetch_coverage_github.py --compare develop + branchCompare mapping. "
            "§4 build_qa_ownership_fields() — qaScope: none skips dev-covered AC from QA TC list. "
            "build_coverage_report.py fills CI {{CI_*}} and Dev tests columns automatically."
        )
        cb.text_frame.paragraphs[0].font.size = Pt(12)
        cb.text_frame.paragraphs[0].font.color.rgb = BODY
        cb.text_frame.paragraphs[0].font.name = FONT
        verdicts = [("Pass", GREEN_BG, GREEN), ("Pass with gaps", AMBER_BG, AMBER), ("Fail", RED_BG, RED)]
        for i, (v, bg, fc) in enumerate(verdicts):
            left = Inches(0.55) + Inches(4.15) * i
            self._rect(s, left, Inches(5.65), Inches(3.95), Inches(0.7), bg, v, 16, True, fc, True)

    def architecture_slide(self):
        s = self.blank()
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHITE
        self.footer(s)
        self._slide_title(s, "Solution architecture", "Cursor IDE · Atlassian MCP · GitHub CLI · Python scripts")
        # center hub
        hub = self._rect(s, Inches(5.0), Inches(2.85), Inches(3.3), Inches(1.5), CORAL, "Coverage\nValidator", 14, True, WHITE, True)
        hub.line.fill.background()
        nodes = [
            (Inches(0.7), Inches(2.5), "Jira MCP", "getJiraIssue\nremote links\nLADR detection"),
            (Inches(10.2), Inches(2.5), "GitHub", "prefetch PR\nchecks · diff"),
            (Inches(0.7), Inches(5.0), "Confluence", "getConfluencePage\nfetch_confluence_\nrequirements.py"),
            (Inches(10.2), Inches(5.0), "HTML report", "reports/\n{KEY}-…-IST.html"),
        ]
        for l, t, title, sub in nodes:
            self._rect(s, l, t, Inches(2.5), Inches(1.35), LIGHT_GRAY, title, 12, True, NAVY, True)
            sb = s.shapes.add_textbox(l + Inches(0.1), t + Inches(0.55), Inches(2.3), Inches(0.7))
            sb.text_frame.text = sub
            sb.text_frame.paragraphs[0].font.size = Pt(9)
            sb.text_frame.paragraphs[0].font.color.rgb = BODY
            sb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            sb.text_frame.paragraphs[0].font.name = FONT
            # connector to hub center approx
            hx, hy = Inches(6.65), Inches(3.55)
            cx = l + Inches(1.25)
            cy = t + Inches(0.68)
            self._arrow(s, cx, cy, hx, hy)
        leg = s.shapes.add_textbox(Inches(0.55), Inches(6.45), Inches(12), Inches(0.45))
        leg.text_frame.text = (
            "Skill: .cursor/skills/msc-dev-code-and-qa-test-coverage-validator/  ·  "
            "fetch_jira_testplan.py merges Confluence LADR + Excel attachment  ·  "
            f"Command: /msc-dev-code-and-qa-test-coverage-validator  ·  Developed by {REPORT_DEVELOPER}"
        )
        leg.text_frame.paragraphs[0].font.size = Pt(10)
        leg.text_frame.paragraphs[0].font.color.rgb = MUTED
        leg.text_frame.paragraphs[0].font.name = FONT

    def report_sections_slide(self):
        s = self.blank()
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHITE
        self.footer(s)
        self._slide_title(
            s,
            "HTML report — 8 sections",
            "reports/{KEY}-{date}-{time}-{TZ}.html  ·  quick links (Jira · plan · PRs · LADR only)  ·  browser-ready",
        )
        sections = [
            ("1", "Coverage summary", "8 metric cards + i tooltips"),
            ("2", "Linked PR(s)", "PR · Repo · State · Title · Files · Dev tests · CI"),
            ("3", "Test plan", "Excel attachment · GWT · Mascot or IDs"),
            ("4", "Dev vs QA", "Ownership split"),
            ("5", "Traceability", "Row-per-AC matrix"),
            ("6", "Impl. review", "Gaps · strengths"),
            ("7", "Assumptions", "Scope notes"),
            ("8", "Actions", "Numbered to-dos"),
        ]
        for i, (num, title, sub) in enumerate(sections):
            row, col = i // 4, i % 4
            left = Inches(0.55) + Inches(3.12) * col
            top = Inches(1.4) + Inches(1.55) * row
            sec_tip = SECTION_HEADER_INFO.get(title, sub)
            card = self._rect(
                s,
                left,
                top,
                Inches(2.95),
                Inches(1.55),
                SOFT_BLUE if row == 0 else SOFT_GOLD,
                radius=True,
            )
            card.line.fill.background()
            n = s.shapes.add_textbox(left + Inches(0.12), top + Inches(0.1), Inches(0.45), Inches(0.4))
            n.text_frame.text = num
            n.text_frame.paragraphs[0].font.size = Pt(20)
            n.text_frame.paragraphs[0].font.bold = True
            n.text_frame.paragraphs[0].font.color.rgb = CORAL
            n.text_frame.paragraphs[0].font.name = FONT
            tt = s.shapes.add_textbox(left + Inches(0.55), top + Inches(0.12), Inches(2.2), Inches(0.4))
            tt.text_frame.text = title
            tt.text_frame.paragraphs[0].font.size = Pt(12)
            tt.text_frame.paragraphs[0].font.bold = True
            tt.text_frame.paragraphs[0].font.color.rgb = NAVY
            tt.text_frame.paragraphs[0].font.name = FONT
            st = s.shapes.add_textbox(left + Inches(0.12), top + Inches(0.75), Inches(2.7), Inches(0.45))
            st.text_frame.text = sub
            st.text_frame.paragraphs[0].font.size = Pt(10)
            st.text_frame.paragraphs[0].font.color.rgb = MUTED
            st.text_frame.paragraphs[0].font.name = FONT
        # mock header bar
        hdr = self._rect(s, Inches(0.55), Inches(4.55), Inches(12.2), Inches(0.55), NAVY, radius=True)
        hdr.line.fill.background()
        ht = s.shapes.add_textbox(Inches(0.75), Inches(4.65), Inches(11.5), Inches(0.38))
        ex = LATEST_EXAMPLE
        ht.text_frame.text = (
            f"{ex['key']}  ·  Verdict: {ex['verdict']}  ·  Generated {ex['generated']}"
        )
        ht.text_frame.paragraphs[0].font.size = Pt(12)
        ht.text_frame.paragraphs[0].font.bold = True
        ht.text_frame.paragraphs[0].font.color.rgb = WHITE
        ht.text_frame.paragraphs[0].font.name = FONT

    def report_ui_slide(self):
        """Report UX — info icons, tooltip layout v8, quick links, Linked PR columns, footer."""
        s = self.blank()
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHITE
        self.footer(s)
        self._slide_title(
            s,
            "Report UX — info icons & tooltips",
            "apply_report_ui_enhancements() — layout v8 · §4 ownership v3 · before every HTML write",
        )
        # Left: where i icons appear
        self._rect(s, Inches(0.55), Inches(1.35), Inches(6.0), Inches(0.42), NAVY, "Where the i icon appears", 12, True, WHITE, True)
        ui_items = [
            "Header verdict + quick links (LADR Confluence filtered)",
            "All 8 section titles + 3 summary group titles",
            "All 8 Coverage summary metric cards + release readiness score",
            "Linked PR(s) — PR, Repo, State, Title, Files, Dev tests, CI status",
            "Test plan + LADR traceability table headers",
            "§4 Dev vs QA — QA execute list excludes dev-covered AC",
        ]
        for j, item in enumerate(ui_items):
            top_i = Inches(1.9) + Inches(0.62) * j
            bx = s.shapes.add_textbox(Inches(0.85), top_i, Inches(5.5), Inches(0.28))
            bx.text_frame.text = item
            bx.text_frame.paragraphs[0].font.size = Pt(10)
            bx.text_frame.paragraphs[0].font.color.rgb = BODY
            bx.text_frame.paragraphs[0].font.name = FONT
            bx.text_frame.word_wrap = True
        # Right: tooltip v8 + quick links
        self._rect(s, Inches(6.85), Inches(1.35), Inches(5.9), Inches(0.42), CORAL, "Tooltip layout v8 · quick links", 12, True, WHITE, True)
        tips = [
            "Quick links: Jira · SharePoint test plan · PR(s) · LADR wiki only",
            "Grooming / deployment / PVC go-live pages excluded from nav",
            "overflow: visible on sections, tables, and headers — no clipped text",
            "PR table: Dev tests + CI status columns with anchored tooltips",
            "Pointer cursor on hover; §5 traceability v2 row evidence lists",
        ]
        for j, tip in enumerate(tips):
            top_t = Inches(1.9) + Inches(0.62) * j
            bx = s.shapes.add_textbox(Inches(7.15), top_t, Inches(5.4), Inches(0.28))
            bx.text_frame.text = tip
            bx.text_frame.paragraphs[0].font.size = Pt(10)
            bx.text_frame.paragraphs[0].font.color.rgb = BODY
            bx.text_frame.paragraphs[0].font.name = FONT
            bx.text_frame.word_wrap = True
        # Linked PR table mock
        self._rect(s, Inches(0.55), Inches(4.85), Inches(12.2), Inches(0.38), NAVY_LIGHT, "§2 Linked PR(s) — six columns", 11, True, WHITE)
        pr_cols = ["PR", "Repo", "State", "Title", "Dev tests", "CI status"]
        cw = Inches(12.2) / 6
        for ci, col in enumerate(pr_cols):
            left = Inches(0.55) + cw * ci
            accent = SOFT_GOLD if col in ("Dev tests", "CI status") else LIGHT_GRAY
            self._rect(s, left, Inches(5.23), cw - Inches(0.04), Inches(0.35), accent, f"{col} i", 9, True, NAVY)
        ex = LATEST_EXAMPLE
        row_txt = f"{ex['pr_note']}  ·  {ex['verdict']}  ·  dev tests {ex['dev_tests_pct']}  ·  CI {ex['ci_line_pct']}"
        self._rect(s, Inches(0.55), Inches(5.58), Inches(12.2), Inches(0.35), WHITE, row_txt, 9, False, BODY)
        # Footer mock
        ft = self._rect(s, Inches(0.55), Inches(6.15), Inches(12.2), Inches(0.45), LIGHT_GRAY, radius=True)
        ft.line.fill.background()
        fb = s.shapes.add_textbox(Inches(0.75), Inches(6.22), Inches(11.8), Inches(0.35))
        fb.text_frame.text = (
            f"Generated by msc-dev-code-and-qa-test-coverage-validator · Developed by {REPORT_DEVELOPER}  "
            f"({ex['report_file']})"
        )
        fb.text_frame.paragraphs[0].font.size = Pt(10)
        fb.text_frame.paragraphs[0].font.color.rgb = MUTED
        fb.text_frame.paragraphs[0].font.name = FONT
        fb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def metrics_dashboard_slide(self):
        s = self.blank()
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHITE
        self.footer(s)
        self._slide_title(
            s,
            "§1 Coverage summary — 8 metric cards",
            f"{PRIMARY_EXAMPLE_KEY} · Pass with gaps · 81% release readiness · thresholds: Green ≥85% · Amber 70–84.9% · Red <70%",
        )
        ex = LATEST_EXAMPLE
        groups = [
            ("Implementation & tests", [
                ("Dev code", ex["dev_code_pct"], GREEN_BG, GREEN),
                ("Dev tests", ex["dev_tests_pct"], AMBER_BG, AMBER),
                ("Req mapped", ex.get("req_mapped", "4/4 AC"), GREEN_BG, GREEN),
            ]),
            ("QA & release risk", [
                ("Test plan AC", ex["testplan_ac_pct"], AMBER_BG, AMBER),
                ("QA remaining", ex["qa_remaining"], LIGHT_GRAY, DARK),
                ("Open gaps", ex["open_gaps"], AMBER_BG, AMBER),
            ]),
            ("CI pipeline", [
                ("Line cov.", ex["ci_line_pct"], GREEN_BG, GREEN),
                ("Branch cov.", ex["ci_branch_pct"], GREEN_BG, GREEN),
            ]),
        ]
        top = Inches(1.35)
        metric_label_map = {
            "Dev code": "Dev code coverage",
            "Dev tests": "Dev unit / integration test coverage",
            "Req mapped": "Requirements mapped",
            "Test plan AC": "Test plan acceptance criteria coverage",
            "QA remaining": "QA scope remaining",
            "Open gaps": "Open gaps",
            "Line cov.": "CI line coverage",
            "Branch cov.": "CI branch coverage",
        }
        for gi, (gname, cards) in enumerate(groups):
            left_base = Inches(0.55) + Inches(4.2) * gi
            gtip = SUMMARY_GROUP_INFO.get(gname, gname)
            gh = s.shapes.add_textbox(left_base, top, Inches(3.5), Inches(0.35))
            gh.text_frame.text = gname
            gh.text_frame.paragraphs[0].font.size = Pt(12)
            gh.text_frame.paragraphs[0].font.bold = True
            gh.text_frame.paragraphs[0].font.color.rgb = NAVY
            gh.text_frame.paragraphs[0].font.name = FONT
            for ci, (lbl, val, bg, vc) in enumerate(cards):
                cy = top + Inches(0.45) + Inches(1.05) * ci
                tip_key = metric_label_map.get(lbl, lbl)
                self._kpi_card(
                    s,
                    left_base,
                    cy,
                    Inches(3.75),
                    Inches(0.95),
                    val,
                    lbl,
                    None,
                    bg,
                    vc,
                )
        # legend
        thresh_tips = {
            "≥85%": "Metric at or above 85% — green in HTML report.",
            "70–84%": "Metric between 70% and 84.9% — amber in HTML report.",
            "<70%": "Metric below 70% — red in HTML report.",
            "NA": "Not available (e.g. no PR for CI, branch-only compare).",
        }
        for i, (lbl, bg, fc) in enumerate([("≥85%", GREEN_BG, GREEN), ("70–84%", AMBER_BG, AMBER), ("<70%", RED_BG, RED), ("NA", LIGHT_GRAY, MUTED)]):
            left = Inches(0.55) + Inches(1.55) * i
            self._rect(
                s,
                left,
                Inches(6.35),
                Inches(1.35),
                Inches(0.48),
                bg,
                lbl,
                10,
                True,
                fc,
                True,
            )

    def testplan_slide(self):
        s = self.blank()
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHITE
        self.footer(s)
        self._slide_title(s, "§3 Test plan validation — Jira Excel + LADR alignment", "Confluence LADR requirements when linked · Excel attachment · Given/When/Then · Evidence")
        feats = [
            ("Confluence LADR (from Jira)", "If story comments or description reference LADR or wiki URLs, agent fetches Confluence via MCP or fetch_confluence_requirements.py", "LADR requirements L1…Ln · cache {KEY}-confluence.json"),
            ("LADR ↔ test plan traceability", "Each L1…Ln requirement tied to Excel test case IDs in report §3", "4/5 LADR mapped · MSC-205625"),
            ("Domino Excel test plan (Jira attachment)", "Jira attachment or local testplans/", "Domino Test Plan.xlsx · Inc as Fulll"),
            ("Given / When / Then", "Content-based Given/When/Then scoring (not column-name only)", "5/5 GWT · MSC-205625"),
            ("LADR + Jira AC mapping", "Semantic match of test cases to Jira AC and Confluence LADR requirements", "77.8% test plan AC · MSC-205625"),
            ("Evidence (Mascot or IDs)", "Mascot URLs when present; else Edit ID from SIT validation (PFT Clear incremental-as-full)", "Edit ID 37ea180e · MSC-205625"),
        ]
        for i, (title, desc, proof) in enumerate(feats):
            top = Inches(1.25) + Inches(0.88) * i
            self._rect(s, Inches(0.55), top, Inches(0.08), Inches(0.95), CORAL)
            tt = s.shapes.add_textbox(Inches(0.75), top, Inches(5.5), Inches(0.35))
            tt.text_frame.text = title
            tt.text_frame.paragraphs[0].font.size = Pt(14)
            tt.text_frame.paragraphs[0].font.bold = True
            tt.text_frame.paragraphs[0].font.color.rgb = NAVY
            tt.text_frame.paragraphs[0].font.name = FONT
            ds = s.shapes.add_textbox(Inches(0.75), top + Inches(0.38), Inches(5.8), Inches(0.5))
            ds.text_frame.text = desc
            ds.text_frame.paragraphs[0].font.size = Pt(11)
            ds.text_frame.paragraphs[0].font.color.rgb = BODY
            ds.text_frame.paragraphs[0].font.name = FONT
            pr = self._rect(s, Inches(6.8), top + Inches(0.15), Inches(5.95), Inches(0.75), SOFT_BLUE, proof, 11, True, NAVY, True)
            pr.line.fill.background()
        # sample table mock
        self._rect(s, Inches(6.8), Inches(5.85), Inches(5.95), Inches(0.4), NAVY, "Scenario examples", 10, True, WHITE)
        for j, row in enumerate(
            [
                "MSC-205625 · Passport LADR · 77.8% test plan AC · 81% readiness",
                "MSC-204417 · Captions LADR · 55.6% test plan AC",
                "MSC-195138 · FF Race · 66.7% test plan AC",
            ]
        ):
            self._rect(s, Inches(6.8), Inches(6.28) + Inches(0.38) * j, Inches(5.95), Inches(0.35), LIGHT_GRAY if j % 2 == 0 else SOFT_BLUE, row, 9, False, BODY)

    def confluence_ladr_slide(self):
        """When Jira references LADR, agent fetches Confluence for supplemental requirements."""
        s = self.blank()
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHITE
        self.footer(s)
        self._slide_title(
            s,
            "Confluence LADR requirements from Jira",
            "When a LADR doc is linked on the story, the agent fetches Confluence and merges LADR requirements with Jira acceptance criteria. Quick nav uses collect_ladr_page_links() — grooming/deployment wiki links stay out of the header.",
        )
        # Trigger column
        self._rect(s, Inches(0.55), Inches(1.35), Inches(3.85), Inches(0.42), NAVY, "Detected from Jira", 12, True, WHITE, True)
        triggers = [
            "Comment or description mentions LADR",
            "Confluence wiki URL (atlassian.net/wiki/…/pages/{id})",
            "Linked design doc referenced in dev/QA comments",
            "Test plan scenarios trace to LADR scope",
        ]
        for j, t in enumerate(triggers):
            bx = s.shapes.add_textbox(Inches(0.75), Inches(1.9) + Inches(0.52) * j, Inches(3.5), Inches(0.45))
            bx.text_frame.text = f"• {t}"
            bx.text_frame.paragraphs[0].font.size = Pt(11)
            bx.text_frame.paragraphs[0].font.color.rgb = BODY
            bx.text_frame.paragraphs[0].font.name = FONT
        # Fetch column
        self._rect(s, Inches(4.65), Inches(1.35), Inches(3.85), Inches(0.42), CORAL, "Agent fetch (one batch)", 12, True, WHITE, True)
        fetch_steps = [
            "getConfluencePage (Atlassian MCP)",
            "python scripts/fetch_confluence_requirements.py {KEY}",
            "Parse LADR requirement items from page body",
            "Write reports/.cache/{KEY}-confluence.json",
        ]
        for j, t in enumerate(fetch_steps):
            bx = s.shapes.add_textbox(Inches(4.85), Inches(1.9) + Inches(0.52) * j, Inches(3.5), Inches(0.45))
            bx.text_frame.text = f"{j + 1}. {t}"
            bx.text_frame.paragraphs[0].font.size = Pt(10)
            bx.text_frame.paragraphs[0].font.color.rgb = BODY
            bx.text_frame.paragraphs[0].font.name = FONT
        # Merge column
        self._rect(s, Inches(8.75), Inches(1.35), Inches(4.0), Inches(0.42), TEAL, "Merged requirement model", 12, True, WHITE, True)
        merge_items = [
            "Jira AC → R1, R2, R3…",
            "Confluence LADR → L1…Ln",
            "Test cases map to R* and L* (semantic match)",
            "Report §3: LADR ↔ test case traceability table",
            "Test plan AC coverage % uses combined set",
        ]
        for j, t in enumerate(merge_items):
            bx = s.shapes.add_textbox(Inches(8.95), Inches(1.9) + Inches(0.52) * j, Inches(3.6), Inches(0.45))
            bx.text_frame.text = f"→ {t}"
            bx.text_frame.paragraphs[0].font.size = Pt(10)
            bx.text_frame.paragraphs[0].font.color.rgb = BODY
            bx.text_frame.paragraphs[0].font.name = FONT
        # Example callout
        self._rect(s, Inches(0.55), Inches(4.15), Inches(12.2), Inches(1.55), SOFT_GOLD, radius=True)
        ex = REPORT_MATRIX[PRIMARY_EXAMPLE_KEY]
        ex_box = s.shapes.add_textbox(Inches(0.75), Inches(4.3), Inches(11.8), Inches(1.25))
        ex_box.text_frame.word_wrap = True
        ex_box.text_frame.text = (
            f"Example — {PRIMARY_EXAMPLE_KEY}: Jira references Passport getting dropped LADR; agent fetches Confluence "
            f"and merges with Domino incremental-as-full test plan ({ex['testplan_note']}). "
            f"Report {ex['report_file']}: {ex['dev_code_pct']} dev code · {ex['dev_tests_pct']} dev tests · "
            f"{ex['testplan_ac_pct']} test plan AC · CI {ex['ci_line_pct']} · {ex['verdict']}. "
            f"Quick links: Jira · SharePoint plan · PR #161 pick-genie · PR #195 encode-monitor · LADR wiki only."
        )
        ex_box.text_frame.paragraphs[0].font.size = Pt(11)
        ex_box.text_frame.paragraphs[0].font.color.rgb = BODY
        ex_box.text_frame.paragraphs[0].font.name = FONT
        # Generic capability strip
        self._rect(s, Inches(0.55), Inches(5.85), Inches(12.2), Inches(0.38), NAVY_LIGHT, "Any LADR-linked story — design requirements become testable L* items for coverage scoring", 10, True, WHITE)
        caps = [
            "Wiki URL or LADR comment on Jira",
            "Confluence page body → requirement list",
            "Merged with Jira AC + Excel test plan",
            "LADR ↔ test case table + gap list for unmapped L*",
        ]
        cw = Inches(12.2) / 4
        for ci, cap in enumerate(caps):
            left = Inches(0.55) + cw * ci
            self._rect(s, left, Inches(6.28), cw - Inches(0.05), Inches(0.55), LIGHT_GRAY, cap, 8, False, NAVY, True)

    def report_matrix_slide(self):
        """Latest HTML report metrics — four recent MSC validations."""
        s = self.blank()
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHITE
        self.footer(s)
        self._slide_title(
            s,
            "Latest report matrix — 8-card coverage summary",
            f"Featured run: {PRIMARY_EXAMPLE_KEY} ({LATEST_EXAMPLE['generated']}) · matrix refreshed from reports/",
        )
        headers = [
            "Ticket",
            "Dev code",
            "Dev tests",
            "Req mapped",
            "Test plan AC",
            "QA scope",
            "Open gaps",
            "CI line",
            "Verdict",
        ]
        col_w = Inches(12.2) / len(headers)
        top = Inches(1.32)
        matrix_header_tips = {
            "Ticket": "MSC Jira issue key for this validation run.",
            "Dev code": SUMMARY_METRIC_INFO["Dev code coverage"],
            "Dev tests": SUMMARY_METRIC_INFO["Dev unit / integration test coverage"],
            "Req mapped": SUMMARY_METRIC_INFO["Requirements mapped"],
            "Test plan AC": SUMMARY_METRIC_INFO["Test plan acceptance criteria coverage"],
            "QA scope": SUMMARY_METRIC_INFO["QA scope remaining"],
            "Open gaps": SUMMARY_METRIC_INFO["Open gaps"],
            "CI line": SUMMARY_METRIC_INFO["CI line coverage"],
            "Verdict": VERDICT_INFO,
        }
        for ci, h in enumerate(headers):
            left = Inches(0.55) + col_w * ci
            ht = matrix_header_tips.get(h, h)
            self._rect(
                s,
                left,
                top,
                col_w - Inches(0.03),
                Inches(0.55),
                NAVY,
                h,
                8,
                True,
                WHITE,
            )
        top = top + Inches(0.38)
        keys = [PRIMARY_EXAMPLE_KEY, "MSC-212571", "MSC-204417", "MSC-195138"]
        for ri, key in enumerate(keys):
            r = REPORT_MATRIX[key]
            row = [
                key,
                r["dev_code_pct"],
                r["dev_tests_pct"],
                r["req_mapped"],
                r["testplan_ac_pct"],
                r["qa_remaining"],
                r["open_gaps"],
                r["ci_line_pct"],
                r["verdict"].replace("Pass with gaps", "Pass w/ gaps"),
            ]
            y = top + Inches(0.42) + Inches(0.38) * ri
            bg = LIGHT_GRAY if ri % 2 else WHITE
            for ci, cell in enumerate(row):
                left = Inches(0.55) + col_w * ci
                fc = DARK if ci == 0 else (AMBER if "Pass" in cell else BODY)
                if ci in (1, 2, 4) and cell.endswith("%"):
                    fc = GREEN if cell.startswith("100") else AMBER if not cell.startswith("66") else AMBER
                if ci == 7 and cell == "NA":
                    fc = MUTED
                self._rect(s, left, y, col_w - Inches(0.03), Inches(0.36), bg, cell, 7, ci == 0, fc)
        # Detail rows — report files + notes
        self._rect(s, Inches(0.55), Inches(2.95), Inches(12.2), Inches(0.32), CORAL, "Report artifacts (latest)", 10, True, WHITE, True)
        for ri, key in enumerate(keys):
            r = REPORT_MATRIX[key]
            note = f"{key} · {r['report_file']} · {r['generated']} · {r['pr_note']} · {r['testplan_note']}"
            if r.get("ladr_note") and r["ladr_note"] != "—":
                note += f" · {r['ladr_note']}"
            y = Inches(3.32) + Inches(0.32) * ri
            highlight = key == LATEST_EXAMPLE["key"]
            self._rect(
                s,
                Inches(0.55),
                y,
                Inches(12.2),
                Inches(0.3),
                SOFT_GOLD if highlight else LIGHT_GRAY,
                note,
                7,
                False,
                NAVY if highlight else BODY,
            )
        # Mini 8-card legend — newest report + branch-only example
        self._rect(s, Inches(0.55), Inches(4.55), Inches(5.9), Inches(0.32), NAVY, f"{LATEST_EXAMPLE['key']} — newest §1 cards", 9, True, WHITE)
        r_new = REPORT_MATRIX[LATEST_EXAMPLE["key"]]
        cards_new = [
            ("Dev code", r_new["dev_code_pct"]),
            ("Dev tests", r_new["dev_tests_pct"]),
            ("Test plan AC", r_new["testplan_ac_pct"]),
            ("CI line", r_new["ci_line_pct"]),
        ]
        for ci, (lbl, val) in enumerate(cards_new):
            left = Inches(0.55) + Inches(1.48) * ci
            fill = SOFT_GOLD if val == "NA" else SOFT_BLUE
            self._kpi_card(s, left, Inches(4.92), Inches(1.38), Inches(0.82), val, lbl, None, fill, NAVY)
        self._rect(s, Inches(6.65), Inches(4.55), Inches(6.1), Inches(0.32), NAVY, "MSC-204417 — branch-only (no PR)", 9, True, WHITE)
        r417 = REPORT_MATRIX["MSC-204417"]
        cards417 = [
            ("Dev code", r417["dev_code_pct"]),
            ("Dev tests", r417["dev_tests_pct"]),
            ("Test plan AC", r417["testplan_ac_pct"]),
            ("CI line", r417["ci_line_pct"]),
        ]
        for ci, (lbl, val) in enumerate(cards417):
            left = Inches(6.65) + Inches(1.48) * ci
            self._kpi_card(s, left, Inches(4.92), Inches(1.38), Inches(0.82), val, lbl, None, SOFT_GOLD if val == "NA" else SOFT_BLUE, NAVY)
        self._rect(
            s,
            Inches(0.55),
            Inches(5.85),
            Inches(12.2),
            Inches(0.32),
            LIGHT_GRAY,
            f"{PRIMARY_EXAMPLE_KEY} — dual PR · passport LADR · 81% release readiness · CI 95%+",
            8,
            False,
            NAVY,
            True,
        )
        r_feat = REPORT_MATRIX[PRIMARY_EXAMPLE_KEY]
        cards_feat = [
            ("Readiness", "81%"),
            ("QA scope", r_feat["qa_remaining"]),
            ("Open gaps", r_feat["open_gaps"]),
            ("Verdict", r_feat["verdict"].replace("Pass with gaps", "Pass w/ gaps")),
        ]
        for ci, (lbl, val) in enumerate(cards_feat):
            left = Inches(0.55) + Inches(1.48) * ci
            self._kpi_card(s, left, Inches(6.22), Inches(1.38), Inches(0.75), val, lbl, None, GREEN_BG, GREEN)

    def dev_qa_slide(self):
        s = self.blank()
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHITE
        self.footer(s)
        self._slide_title(s, "§4–§5 Dev vs QA ownership & traceability", "No ambiguity on what QA must still run in SIT/UAT")
        self._rect(s, Inches(0.55), Inches(1.4), Inches(5.9), Inches(4.9), SOFT_BLUE, radius=True)
        dh = s.shapes.add_textbox(Inches(0.75), Inches(1.55), Inches(5.5), Inches(0.4))
        dh.text_frame.text = "Development proof"
        dh.text_frame.paragraphs[0].font.size = Pt(16)
        dh.text_frame.paragraphs[0].font.bold = True
        dh.text_frame.paragraphs[0].font.color.rgb = NAVY
        dh.text_frame.paragraphs[0].font.name = FONT
        dev_items = [
            "Unit / integration tests in PR",
            "File + test name evidence",
            "Dev-owned AC marked Covered",
            f"CI line coverage — {LATEST_EXAMPLE['ci_line_pct']} (linked PR)",
            "Dev-covered AC → qaScope: none (excluded from QA TC execute list)",
        ]
        for j, it in enumerate(dev_items):
            bx = s.shapes.add_textbox(Inches(0.85), Inches(2.1) + Inches(0.55) * j, Inches(5.3), Inches(0.45))
            bx.text_frame.text = f"✓  {it}"
            bx.text_frame.paragraphs[0].font.size = Pt(12)
            bx.text_frame.paragraphs[0].font.color.rgb = GREEN
            bx.text_frame.paragraphs[0].font.name = FONT
        self._rect(s, Inches(6.85), Inches(1.4), Inches(5.9), Inches(4.9), SOFT_CORAL, radius=True)
        qh = s.shapes.add_textbox(Inches(7.05), Inches(1.55), Inches(5.5), Inches(0.4))
        qh.text_frame.text = "QA handoff"
        qh.text_frame.paragraphs[0].font.size = Pt(16)
        qh.text_frame.paragraphs[0].font.bold = True
        qh.text_frame.paragraphs[0].font.color.rgb = CORAL
        qh.text_frame.paragraphs[0].font.name = FONT
        qa_items = [
            "E2E / manual for AC without dev test proof",
            "Execute only TCs mapped to QA-scoped requirements (§4 filter)",
            "Mascot or SIT Jobs IDs in test plan Evidence column",
            "Spot-check LADR status codes when unit tests partial",
        ]
        for j, it in enumerate(qa_items):
            bx = s.shapes.add_textbox(Inches(7.15), Inches(2.1) + Inches(0.55) * j, Inches(5.3), Inches(0.45))
            bx.text_frame.text = f"→  {it}"
            bx.text_frame.paragraphs[0].font.size = Pt(12)
            bx.text_frame.paragraphs[0].font.color.rgb = CORAL
            bx.text_frame.paragraphs[0].font.name = FONT
        mid = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.15), Inches(3.5), Inches(1.0), Inches(1.0))
        mid.fill.solid()
        mid.fill.fore_color.rgb = GOLD
        mid.line.fill.background()
        mid.text_frame.text = "AC"
        mid.text_frame.paragraphs[0].font.size = Pt(14)
        mid.text_frame.paragraphs[0].font.bold = True
        mid.text_frame.paragraphs[0].font.color.rgb = NAVY
        mid.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        mid.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _matrix_header(self, slide, left, top, width, text, fill=NAVY):
        self._rect(slide, left, top, width, Inches(0.42), fill, text, 10, True, WHITE)

    def _add_bullets(self, slide, left, top, width, height, lines: list[str], *, font_pt: int = 10, color=BODY):
        if not lines:
            return
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines[:14]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(font_pt)
            p.font.name = FONT
            p.font.color.rgb = color
            p.space_after = Pt(3)

    def _cell_status_style(self, text: str, *, default_bg=WHITE) -> tuple:
        t = (text or "").lower()
        if any(x in t for x in ("implemented", "covered", "passed", "merged", "full", "aligns", "complete")):
            return GREEN_BG, GREEN
        if any(x in t for x in ("partial", "gap", "missing", "medium", "concern", "warn")):
            return AMBER_BG, AMBER
        if any(x in t for x in ("fail", "not implemented")):
            return RED_BG, RED
        if "e2e" in t or "manual" in t:
            return SOFT_CORAL, CORAL
        if "%" in text:
            try:
                pct = float(re.search(r"[\d.]+", text).group())
                if pct >= 85:
                    return GREEN_BG, GREEN
                if pct >= 70:
                    return AMBER_BG, AMBER
                return RED_BG, RED
            except (AttributeError, ValueError):
                pass
        return default_bg, DARK if default_bg == WHITE else BODY

    def _section_accent(self, num: str) -> tuple:
        accents = {
            "1": (SOFT_BLUE, NAVY),
            "2": (SOFT_GOLD, GOLD_DARK),
            "3": (SOFT_CORAL, CORAL),
            "4": (RGBColor(0xE0, 0xE7, 0xFF), NAVY),
            "5": (RGBColor(0xCC, 0xFB, 0xF1), TEAL),
            "6": (GREEN_BG, GREEN),
            "7": (LIGHT_GRAY, BODY),
            "8": (RGBColor(0xFE, 0xF3, 0xC7), AMBER),
        }
        return accents.get(num, (SOFT_BLUE, NAVY))

    def _section_badge(self, slide, num: str, left, top, size=Inches(0.55)):
        bg, fg = self._section_accent(num)
        o = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
        o.fill.solid()
        o.fill.fore_color.rgb = CORAL if num == "3" else NAVY
        o.line.fill.background()
        o.text_frame.text = num
        p = o.text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER
        o.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        return bg, fg

    def _kpi_tile(self, slide, value, label, left, top, w, h, fill, value_color, label_color=MUTED):
        card = self._rect(slide, left, top, w, h, fill, radius=True)
        card.line.fill.background()
        v = slide.shapes.add_textbox(left, top + Inches(0.1), w, Inches(0.42))
        v.text_frame.text = value
        v.text_frame.paragraphs[0].font.size = Pt(22)
        v.text_frame.paragraphs[0].font.bold = True
        v.text_frame.paragraphs[0].font.color.rgb = value_color
        v.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        v.text_frame.paragraphs[0].font.name = FONT
        lb = slide.shapes.add_textbox(left, top + Inches(0.5), w, Inches(0.35))
        lb.text_frame.text = label
        lb.text_frame.paragraphs[0].font.size = Pt(8)
        lb.text_frame.paragraphs[0].font.color.rgb = label_color
        lb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        lb.text_frame.paragraphs[0].font.name = FONT

    def _add_styled_table(
        self,
        slide,
        left,
        top,
        width,
        headers: list[str],
        rows: list[list[str]],
        *,
        row_h=Inches(0.38),
        font_pt: int = 8,
    ):
        if not headers or not rows:
            return
        n = len(headers)
        col_w = width / n
        for ci, h in enumerate(headers):
            self._rect(
                slide, left + col_w * ci, top, col_w - Inches(0.03), Inches(0.36),
                NAVY, h[:16], font_pt, True, WHITE, True,
            )
        for ri, row in enumerate(rows[:6]):
            y = top + Inches(0.38) + row_h * ri
            for ci in range(n):
                cell = row[ci] if ci < len(row) else ""
                bg, fg = self._cell_status_style(cell, default_bg=LIGHT_GRAY if ri % 2 else WHITE)
                self._rect(
                    slide, left + col_w * ci, y, col_w - Inches(0.03), row_h - Inches(0.04),
                    bg, cell[:40], font_pt, ci == 0, fg, True,
                )

    def _report_section_header(self, slide, sec: ReportSection, issue_key: str) -> tuple:
        """Navy band + § badge; returns (accent_bg, accent_fg, content_top)."""
        bg_acc, fg_acc = self._section_accent(sec.num)
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.95))
        band.fill.solid()
        band.fill.fore_color.rgb = bg_acc
        band.line.fill.background()
        self._section_badge(slide, sec.num, Inches(0.7), Inches(0.42))
        tit = slide.shapes.add_textbox(Inches(1.35), Inches(0.38), Inches(10.5), Inches(0.4))
        tit.text_frame.text = sec.title
        tit.text_frame.paragraphs[0].font.size = Pt(22)
        tit.text_frame.paragraphs[0].font.bold = True
        tit.text_frame.paragraphs[0].font.color.rgb = fg_acc
        tit.text_frame.paragraphs[0].font.name = FONT
        sub = sec.lead if sec.lead else f"{issue_key} · live HTML report section"
        if len(sub) > 100:
            sub = sub[:97] + "…"
        st = slide.shapes.add_textbox(Inches(1.35), Inches(0.78), Inches(10.8), Inches(0.32))
        st.text_frame.text = sub
        st.text_frame.paragraphs[0].font.size = Pt(10)
        st.text_frame.paragraphs[0].font.color.rgb = BODY
        st.text_frame.paragraphs[0].font.name = FONT
        self.footer(slide)
        return bg_acc, fg_acc, Inches(1.38)

    def report_overview_slide(self, data: dict) -> None:
        """Slide 11 — hero dashboard + 8-section navigator."""
        s = self.blank()
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHITE
        self.footer(s)
        key = data["issue_key"]
        ex = LATEST_EXAMPLE
        hero = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(2.05))
        hero.fill.solid()
        hero.fill.fore_color.rgb = NAVY
        hero.line.fill.background()
        ht = s.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(8.5), Inches(0.55))
        ht.text_frame.text = f"{key} — Live coverage validation report"
        ht.text_frame.paragraphs[0].font.size = Pt(24)
        ht.text_frame.paragraphs[0].font.bold = True
        ht.text_frame.paragraphs[0].font.color.rgb = WHITE
        ht.text_frame.paragraphs[0].font.name = FONT
        story = data["story_title"]
        if len(story) > 68:
            story = story[:65] + "…"
        hs = s.shapes.add_textbox(Inches(0.55), Inches(0.82), Inches(8.8), Inches(0.55))
        hs.text_frame.text = story
        hs.text_frame.word_wrap = True
        hs.text_frame.paragraphs[0].font.size = Pt(11)
        hs.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
        hs.text_frame.paragraphs[0].font.name = FONT
        badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.2), Inches(0.45), Inches(2.55), Inches(1.15))
        badge.fill.solid()
        badge.fill.fore_color.rgb = AMBER_BG
        badge.line.color.rgb = AMBER
        badge.line.width = Pt(1.5)
        badge.text_frame.text = f"{data['verdict']}\n{data['readiness']} ready"
        badge.text_frame.paragraphs[0].font.size = Pt(13)
        badge.text_frame.paragraphs[0].font.bold = True
        badge.text_frame.paragraphs[0].font.color.rgb = AMBER
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge.text_frame.paragraphs[0].font.name = FONT
        kpis = [
            (ex["dev_code_pct"], "Dev code", GREEN_BG, GREEN),
            (ex["dev_tests_pct"], "Dev tests", SOFT_BLUE, NAVY),
            (ex["testplan_ac_pct"], "Test plan AC", AMBER_BG, AMBER),
            (ex["ci_line_pct"], "CI line", GREEN_BG, GREEN),
        ]
        kw = Inches(2.95)
        for i, (val, lbl, fill, vc) in enumerate(kpis):
            self._kpi_tile(s, val, lbl, Inches(0.55) + (kw + Inches(0.12)) * i, Inches(2.2), kw, Inches(1.05), fill, vc)
        self._rect(s, Inches(0.55), Inches(3.4), Inches(12.2), Inches(0.36), CORAL, "Report walkthrough — 8 sections", 12, True, WHITE, True)
        tiles = data.get("sections", [])
        tw, th = Inches(2.9), Inches(0.72)
        for i, sec in enumerate(tiles[:8]):
            col, row = i % 4, i // 4
            left = Inches(0.55) + (tw + Inches(0.15)) * col
            top = Inches(3.88) + (th + Inches(0.12)) * row
            fill, accent = self._section_accent(sec.num)
            self._rect(s, left, top, tw, th, fill, radius=True)
            nb = s.shapes.add_textbox(left + Inches(0.12), top + Inches(0.08), Inches(0.4), Inches(0.28))
            nb.text_frame.text = sec.num
            nb.text_frame.paragraphs[0].font.size = Pt(14)
            nb.text_frame.paragraphs[0].font.bold = True
            nb.text_frame.paragraphs[0].font.color.rgb = CORAL if sec.num == "3" else NAVY
            nb.text_frame.paragraphs[0].font.name = FONT
            tt = s.shapes.add_textbox(left + Inches(0.45), top + Inches(0.06), tw - Inches(0.5), Inches(0.55))
            tt.text_frame.text = sec.title[:28] + ("…" if len(sec.title) > 28 else "")
            tt.text_frame.word_wrap = True
            tt.text_frame.paragraphs[0].font.size = Pt(9)
            tt.text_frame.paragraphs[0].font.bold = True
            tt.text_frame.paragraphs[0].font.color.rgb = accent
            tt.text_frame.paragraphs[0].font.name = FONT
        ready = s.shapes.add_textbox(Inches(0.55), Inches(5.55), Inches(12.2), Inches(1.2))
        ready.text_frame.word_wrap = True
        ready_lines = ["✓ Jira readiness"] + [f"  ✓ {it}" for it in data.get("readiness_items", [])[:4]]
        self._add_bullets(s, Inches(0.55), Inches(5.55), Inches(5.8), Inches(1.15), ready_lines, font_pt=9)
        src = s.shapes.add_textbox(Inches(6.6), Inches(5.55), Inches(6.1), Inches(0.9))
        src.text_frame.word_wrap = True
        src.text_frame.text = (
            f"Source: {data['report_file']}\n"
            f"Generated {data.get('generated', '')}\n"
            f"Open the HTML report for full tables, tooltips, and evidence links."
        )
        for p in src.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.color.rgb = MUTED
            p.font.name = FONT

    _SUMMARY_LABEL_SHORT = {
        "release readiness score": "Release readiness",
        "dev code coverage": "Dev code",
        "dev unit / integration test coverage": "Dev tests",
        "requirements mapped": "Requirements mapped",
        "test plan acceptance criteria coverage": "Test plan AC",
        "qa scope remaining": "QA scope remaining",
        "open gaps": "Open gaps",
        "ci line coverage": "CI line",
        "ci branch coverage": "CI branch",
    }

    _SUMMARY_GROUPS = (
        ("Implementation & tests", ("dev code coverage", "dev unit / integration test coverage", "requirements mapped")),
        ("QA & release risk", ("test plan acceptance criteria coverage", "qa scope remaining", "open gaps")),
        ("CI pipeline", ("ci line coverage", "ci branch coverage")),
    )

    def _parse_summary_metric(self, line: str) -> tuple[str, str, str]:
        """Return (label_key, value, footnote) from 'Label: value — note'."""
        if ":" not in line:
            return "", line.strip(), ""
        label, rest = line.split(":", 1)
        label_key = label.strip().lower()
        if " — " in rest:
            val, note = rest.split(" — ", 1)
        else:
            val, note = rest.strip(), ""
        return label_key, val.strip(), note.strip()[:55]

    def _summary_metric_card(
        self, slide, left, top, width, height, value: str, label: str, footnote: str = "",
    ) -> None:
        fill, vc = self._cell_status_style(value)
        sh = self._rect(slide, left, top, width, height, fill, radius=True)
        if sh is not None:
            sh.line.fill.background()
        v = slide.shapes.add_textbox(left, top + Inches(0.08), width, Inches(0.38))
        v.text_frame.text = value[:18]
        v.text_frame.paragraphs[0].font.size = Pt(20)
        v.text_frame.paragraphs[0].font.bold = True
        v.text_frame.paragraphs[0].font.color.rgb = vc
        v.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        v.text_frame.paragraphs[0].font.name = FONT
        lb = slide.shapes.add_textbox(left + Inches(0.06), top + Inches(0.46), width - Inches(0.12), Inches(0.28))
        lb.text_frame.text = label
        lb.text_frame.paragraphs[0].font.size = Pt(9)
        lb.text_frame.paragraphs[0].font.bold = True
        lb.text_frame.paragraphs[0].font.color.rgb = NAVY
        lb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        lb.text_frame.paragraphs[0].font.name = FONT
        if footnote:
            fn = slide.shapes.add_textbox(left + Inches(0.06), top + height - Inches(0.32), width - Inches(0.12), Inches(0.28))
            fn.text_frame.word_wrap = True
            fn.text_frame.text = footnote
            fn.text_frame.paragraphs[0].font.size = Pt(7)
            fn.text_frame.paragraphs[0].font.color.rgb = MUTED
            fn.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            fn.text_frame.paragraphs[0].font.name = FONT

    def _slide_section_summary(self, s, sec: ReportSection, top) -> None:
        """§1 Coverage summary — hero readiness + grouped metric cards."""
        by_key: dict[str, tuple[str, str]] = {}
        for line in sec.bullets:
            key, val, note = self._parse_summary_metric(line)
            if not key or key in by_key:
                continue
            by_key[key] = (val, note)

        if not by_key:
            self._add_bullets(s, Inches(0.55), top, Inches(12.2), Inches(5.2), sec.bullets, font_pt=10)
            return

        readiness = by_key.get("release readiness score")
        if readiness:
            val, note = readiness
            fill, vc = self._cell_status_style(val)
            hero_h = Inches(1.05)
            self._rect(s, Inches(0.55), top, Inches(12.2), hero_h, fill, radius=True)
            hv = s.shapes.add_textbox(Inches(0.75), top + Inches(0.12), Inches(2.2), Inches(0.55))
            hv.text_frame.text = val
            hv.text_frame.paragraphs[0].font.size = Pt(36)
            hv.text_frame.paragraphs[0].font.bold = True
            hv.text_frame.paragraphs[0].font.color.rgb = vc
            hv.text_frame.paragraphs[0].font.name = FONT
            hl = s.shapes.add_textbox(Inches(2.95), top + Inches(0.18), Inches(4.5), Inches(0.4))
            hl.text_frame.text = "Release readiness score"
            hl.text_frame.paragraphs[0].font.size = Pt(14)
            hl.text_frame.paragraphs[0].font.bold = True
            hl.text_frame.paragraphs[0].font.color.rgb = NAVY
            hl.text_frame.paragraphs[0].font.name = FONT
            if note:
                hn = s.shapes.add_textbox(Inches(2.95), top + Inches(0.52), Inches(8.5), Inches(0.4))
                hn.text_frame.text = note
                hn.text_frame.paragraphs[0].font.size = Pt(9)
                hn.text_frame.paragraphs[0].font.color.rgb = MUTED
                hn.text_frame.paragraphs[0].font.name = FONT
            top = top + hero_h + Inches(0.14)

        col_w = Inches(5.95)
        card_w = Inches(1.85)
        card_h = Inches(1.05)
        gap = Inches(0.1)
        row1_top = top

        for gi, (group_title, keys) in enumerate(self._SUMMARY_GROUPS[:2]):
            metrics = [(k, by_key[k]) for k in keys if k in by_key]
            if not metrics:
                continue
            left = Inches(0.55) + (col_w + Inches(0.3)) * gi
            self._rect(s, left, row1_top, col_w, Inches(0.32), NAVY, group_title, 10, True, WHITE, True)
            row_top = row1_top + Inches(0.38)
            for i, (key, (val, note)) in enumerate(metrics):
                short = self._SUMMARY_LABEL_SHORT.get(key, key.title()[:20])
                self._summary_metric_card(
                    s, left + (card_w + gap) * i, row_top, card_w, card_h, val, short, note,
                )

        ci_keys = self._SUMMARY_GROUPS[2][1]
        ci_metrics = [(k, by_key[k]) for k in ci_keys if k in by_key]
        if ci_metrics:
            ci_top = row1_top + card_h + Inches(0.55)
            gw = Inches(12.2)
            self._rect(s, Inches(0.55), ci_top, gw, Inches(0.32), NAVY, "CI pipeline", 10, True, WHITE, True)
            row_top = ci_top + Inches(0.38)
            n = len(ci_metrics)
            total_w = n * card_w + max(0, n - 1) * gap
            start = Inches(0.55) + (gw - total_w) / 2
            for i, (key, (val, note)) in enumerate(ci_metrics):
                short = self._SUMMARY_LABEL_SHORT.get(key, key.title()[:20])
                self._summary_metric_card(
                    s, start + (card_w + gap) * i, row_top, card_w, card_h, val, short, note,
                )

    def _slide_section_prs(self, s, sec: ReportSection, top) -> None:
        for i, row in enumerate(sec.table_rows[:2]):
            left = Inches(0.55) + Inches(6.15) * i
            merged = "MERGED" in " ".join(row).upper()
            fill = GREEN_BG if merged else SOFT_GOLD
            self._rect(s, left, top, Inches(5.95), Inches(2.35), fill, radius=True)
            pr = row[0] if row else f"PR {i+1}"
            repo = row[1] if len(row) > 1 else ""
            state = row[2] if len(row) > 2 else ""
            title = row[3] if len(row) > 3 else ""
            tests = row[5] if len(row) > 5 else ""
            ci = row[6] if len(row) > 6 else ""
            h = s.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), Inches(5.5), Inches(0.35))
            h.text_frame.text = f"{pr}  ·  {repo}  ·  {state}"
            h.text_frame.paragraphs[0].font.size = Pt(12)
            h.text_frame.paragraphs[0].font.bold = True
            h.text_frame.paragraphs[0].font.color.rgb = NAVY
            h.text_frame.paragraphs[0].font.name = FONT
            body = s.shapes.add_textbox(left + Inches(0.2), top + Inches(0.55), Inches(5.5), Inches(1.0))
            body.text_frame.word_wrap = True
            body.text_frame.text = f"{title[:120]}\n\nDev tests: {tests[:70]}\nCI: {ci}"
            for p in body.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = BODY
                p.font.name = FONT

    def _slide_section_ownership(self, s, sec: ReportSection, top) -> None:
        dev_lines, qa_lines = [], []
        block = "dev"
        for line in sec.bullets:
            if line.startswith("QA handoff"):
                block = "qa"
                continue
            if line.startswith("Dev-owned"):
                continue
            (qa_lines if block == "qa" else dev_lines).append(line)
        self._rect(s, Inches(0.55), top, Inches(5.9), Inches(4.5), SOFT_BLUE, radius=True)
        dh = s.shapes.add_textbox(Inches(0.75), top + Inches(0.12), Inches(5.5), Inches(0.35))
        dh.text_frame.text = "Development proof"
        dh.text_frame.paragraphs[0].font.size = Pt(14)
        dh.text_frame.paragraphs[0].font.bold = True
        dh.text_frame.paragraphs[0].font.color.rgb = NAVY
        dh.text_frame.paragraphs[0].font.name = FONT
        self._add_bullets(s, Inches(0.75), top + Inches(0.5), Inches(5.5), Inches(3.8), dev_lines, font_pt=10, color=GREEN)
        self._rect(s, Inches(6.85), top, Inches(5.9), Inches(4.5), SOFT_CORAL, radius=True)
        qh = s.shapes.add_textbox(Inches(7.05), top + Inches(0.12), Inches(5.5), Inches(0.35))
        qh.text_frame.text = "QA handoff"
        qh.text_frame.paragraphs[0].font.size = Pt(14)
        qh.text_frame.paragraphs[0].font.bold = True
        qh.text_frame.paragraphs[0].font.color.rgb = CORAL
        qh.text_frame.paragraphs[0].font.name = FONT
        self._add_bullets(s, Inches(7.15), top + Inches(0.5), Inches(5.5), Inches(3.8), qa_lines, font_pt=10, color=CORAL)

    def _slide_section_review(self, s, sec: ReportSection, top) -> None:
        pos, gaps = [], []
        mode = "pos"
        for line in sec.bullets:
            if "Gaps" in line or "concerns" in line:
                mode = "gap"
                continue
            if "Correctly" in line:
                continue
            (gaps if mode == "gap" else pos).append(line.lstrip())
        self._rect(s, Inches(0.55), top, Inches(5.95), Inches(2.0), GREEN_BG, radius=True)
        self._add_bullets(s, Inches(0.75), top + Inches(0.15), Inches(5.5), Inches(1.65), ["✓ Strengths"] + pos, font_pt=11, color=GREEN)
        self._rect(s, Inches(6.65), top, Inches(5.95), Inches(2.0), RED_BG, radius=True)
        self._add_bullets(s, Inches(6.85), top + Inches(0.15), Inches(5.5), Inches(1.65), ["⚠ Gaps"] + gaps, font_pt=11, color=RED)

    def _slide_section_actions(self, s, sec: ReportSection, top) -> None:
        for i, action in enumerate(sec.bullets[:5]):
            y = top + Inches(0.85) * i
            self._rect(s, Inches(0.55), y, Inches(12.2), Inches(0.72), SOFT_GOLD if i % 2 else LIGHT_GRAY, radius=True)
            num = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.14), Inches(0.42), Inches(0.42))
            num.fill.solid()
            num.fill.fore_color.rgb = GOLD
            num.line.fill.background()
            num.text_frame.text = str(i + 1)
            num.text_frame.paragraphs[0].font.size = Pt(12)
            num.text_frame.paragraphs[0].font.bold = True
            num.text_frame.paragraphs[0].font.color.rgb = NAVY
            num.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            num.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            tx = s.shapes.add_textbox(Inches(1.25), y + Inches(0.12), Inches(11.2), Inches(0.5))
            tx.text_frame.text = action
            tx.text_frame.paragraphs[0].font.size = Pt(12)
            tx.text_frame.paragraphs[0].font.color.rgb = NAVY
            tx.text_frame.paragraphs[0].font.name = FONT

    def report_section_slide(self, sec: ReportSection, *, issue_key: str = PRIMARY_EXAMPLE_KEY) -> None:
        """One styled slide per HTML report section (§1–§8)."""
        s = self.blank()
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHITE
        _, _, top = self._report_section_header(s, sec, issue_key)
        if sec.num == "1":
            self._slide_section_summary(s, sec, top)
        elif sec.num == "2" and sec.table_rows:
            self._slide_section_prs(s, sec, top)
        elif sec.num == "4":
            self._slide_section_ownership(s, sec, top)
        elif sec.num == "6":
            self._slide_section_review(s, sec, top)
        elif sec.num == "8":
            self._slide_section_actions(s, sec, top)
        else:
            if sec.bullets:
                callout = self._rect(s, Inches(0.55), top, Inches(12.2), Inches(0.95), SOFT_BLUE, radius=True)
                callout.line.fill.background()
                self._add_bullets(s, Inches(0.75), top + Inches(0.12), Inches(11.6), Inches(0.75), sec.bullets[:4], font_pt=10)
                top = top + Inches(1.08)
            if sec.table_rows and sec.table_headers:
                self._add_styled_table(s, Inches(0.55), top, Inches(12.2), sec.table_headers, sec.table_rows)
            elif sec.bullets and sec.num not in ("1", "2", "4", "6", "8"):
                self._add_bullets(s, Inches(0.55), top, Inches(12.2), Inches(5.2), sec.bullets, font_pt=10)

    def build_report_tail_from_html(self, data: dict) -> None:
        """Slides 11+ — enablement, dashboard, 8 HTML sections, proven outcomes, closing."""
        self.setup_slide()
        self.report_overview_slide(data)
        for sec in data["sections"]:
            self.report_section_slide(sec, issue_key=data["issue_key"])
        self.case_studies_slide()
        self.closing_slide()

    def _matrix_cell(
        self,
        slide,
        left,
        top,
        width,
        height,
        lines: list[str],
        fill=WHITE,
        title: str | None = None,
        title_color=NAVY,
        line_colors: list[RGBColor] | None = None,
    ):
        self._rect(slide, left, top, width, height, fill, radius=False)
        box = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.06), width - Inches(0.16), height - Inches(0.1))
        tf = box.text_frame
        tf.word_wrap = True
        if title:
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = title_color
            p.font.name = FONT
        for i, line in enumerate(lines):
            para = tf.add_paragraph() if title or i > 0 else tf.paragraphs[0]
            para.text = line
            para.font.size = Pt(9)
            para.font.name = FONT
            para.font.color.rgb = (line_colors[i] if line_colors and i < len(line_colors) else BODY)
            para.space_after = Pt(2)

    def requirement_traceability_slide(self):
        """MSC-205625 example — Jira requirements vs PR code coverage vs Jira-attached Excel test plan."""
        s = self.blank()
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHITE
        self.footer(s)
        self._slide_title(
            s,
            "Requirement traceability chart — Jira · Code · Test plan",
            f"{LATEST_EXAMPLE['key']} · {LATEST_EXAMPLE['testplan_note']} · {LATEST_EXAMPLE['pr_note']}",
        )
        left = Inches(0.55)
        col_req = Inches(4.35)
        col_code = Inches(4.05)
        col_plan = Inches(3.95)
        top = Inches(1.35)
        self._matrix_header(s, left, top, col_req, "Requirement (Jira story)")
        self._matrix_header(s, left + col_req + Inches(0.08), top, col_code, "Code coverage (linked PR)")
        self._matrix_header(s, left + col_req + col_code + Inches(0.16), top, col_plan, "Test plan (Jira Excel attachment)")
        rows = [
            (
                "R1",
                "Content passport retained in cumulative output manifestation for PFT Clear incremental-as-full",
                ["Implemented", "Dev tests: Covered", "PR #161 · passport_manager.py"],
                [GREEN, GREEN, BODY],
                GREEN_BG,
                ["TC1–TC5 mapped", "Given/When/Then: full", "Mascot evidence linked"],
                [GREEN, GREEN, GREEN],
                GREEN_BG,
            ),
            (
                "R2",
                "When metadata-update becomes full fulfillment, pack passport must be delivered for incremental-as-full",
                ["Implemented", "Dev tests: Partial", "PriorFulfilledManifestRecord model"],
                [GREEN, AMBER, BODY],
                AMBER_BG,
                ["TC1, TC4 mapped", "Given/When/Then: full", "Aligns with PR behavior"],
                [GREEN, GREEN, GREEN],
                SOFT_BLUE,
            ),
            (
                "R3",
                "Pick-genie must not drop passport re-fetch when workflow resolves fulfillmentType full (PFT Clear)",
                ["Implemented", "Dev tests: Covered", "TestDominoPassportRouting"],
                [GREEN, GREEN, BODY],
                GREEN_BG,
                ["TC1–TC5 mapped", "Given/When/Then: full", "Incremental-as-full scenarios"],
                [GREEN, GREEN, GREEN],
                GREEN_BG,
            ),
            (
                "R4",
                "Fix validated in SIT with Edit ID 37ea180e-77cc-413f-95cf-9dfcebf08cd2 and provided Mascot evidence",
                ["Partial", "Dev tests: Missing", "QA Manual — SIT validation"],
                [AMBER, RED, BODY],
                AMBER_BG,
                ["No mapped test case", "Gap: missing Edit ID scenario", "Then steps lack passport assertion"],
                [RED, RED, AMBER],
                RED_BG,
            ),
        ]
        row_h = Inches(1.05)
        for ri, (rid, req_text, code_lines, code_colors, code_bg, plan_lines, plan_colors, plan_bg) in enumerate(rows):
            y = top + Inches(0.48) + row_h * ri
            req_title = f"{rid}"
            req_body = req_text if len(req_text) <= 95 else req_text[:92] + "…"
            self._matrix_cell(s, left, y, col_req, row_h, [req_body], fill=LIGHT_GRAY if ri % 2 else WHITE, title=req_title)
            self._matrix_cell(
                s,
                left + col_req + Inches(0.08),
                y,
                col_code,
                row_h,
                code_lines,
                fill=code_bg,
                line_colors=code_colors,
            )
            self._matrix_cell(
                s,
                left + col_req + col_code + Inches(0.16),
                y,
                col_plan,
                row_h,
                plan_lines,
                fill=plan_bg,
                line_colors=plan_colors,
            )
        note = s.shapes.add_textbox(Inches(0.55), Inches(6.15), Inches(12.2), Inches(0.55))
        note.text_frame.word_wrap = True
        note.text_frame.text = (
            "Left: acceptance criteria extracted from Jira description · Middle: production code + dev unit/integration tests from PR diff · "
            "Right: scenarios parsed from the Excel test plan file attached to the Jira ticket (not a separate tool — the workbook in Jira attachments)."
        )
        note.text_frame.paragraphs[0].font.size = Pt(10)
        note.text_frame.paragraphs[0].font.color.rgb = MUTED
        note.text_frame.paragraphs[0].font.name = FONT

    def setup_slide(self):
        s = self.blank()
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHITE
        self.footer(s)
        self._slide_title(s, "Enablement — one-time setup (~10 min)", "Clone lab → MCP → gh → permissions → optional .env")
        steps = [
            ("Clone pegasus-qa-agents-lab", "cursor ."),
            ("pip install -r requirements.txt", "Python 3.10+"),
            ("Atlassian MCP auth", "wbdstreaming.atlassian.net"),
            ("gh auth login", "GitHub PR access"),
            ("install_coverage_validator_permissions.py", "Auto-run allowlist"),
            (".env + validator.defaults.json", "Test plan download + repo defaults"),
        ]
        for i, (title, sub) in enumerate(steps):
            row = i // 2
            col = i % 2
            left = Inches(0.55) + Inches(6.2) * col
            top = Inches(1.35) + Inches(1.05) * row
            self._rect(s, left, top, Inches(5.95), Inches(0.88), LIGHT_GRAY if i % 2 else SOFT_GOLD, radius=True)
            check = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.15), top + Inches(0.22), Inches(0.45), Inches(0.45))
            check.fill.solid()
            check.fill.fore_color.rgb = GREEN
            check.line.fill.background()
            check.text_frame.text = "✓"
            check.text_frame.paragraphs[0].font.size = Pt(14)
            check.text_frame.paragraphs[0].font.bold = True
            check.text_frame.paragraphs[0].font.color.rgb = WHITE
            check.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            check.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            tt = s.shapes.add_textbox(left + Inches(0.75), top + Inches(0.12), Inches(5.0), Inches(0.35))
            tt.text_frame.text = title
            tt.text_frame.paragraphs[0].font.size = Pt(12)
            tt.text_frame.paragraphs[0].font.bold = True
            tt.text_frame.paragraphs[0].font.color.rgb = NAVY
            tt.text_frame.paragraphs[0].font.name = FONT
            st = s.shapes.add_textbox(left + Inches(0.75), top + Inches(0.48), Inches(5.0), Inches(0.3))
            st.text_frame.text = sub
            st.text_frame.paragraphs[0].font.size = Pt(10)
            st.text_frame.paragraphs[0].font.color.rgb = MUTED
            st.text_frame.paragraphs[0].font.name = FONT
        scripts = s.shapes.add_textbox(Inches(0.55), Inches(4.55), Inches(12.2), Inches(1.8))
        scripts.text_frame.word_wrap = True
        lines = [
            "Scripts (batched — no manual gh per run):",
            "fetch_confluence_requirements.py  ·  fetch_jira_testplan.py  ·  prefetch_coverage_inputs.py",
            "map_requirements_to_diff.py  ·  build_coverage_report.py  ·  apply_report_ui_enhancements()",
            f"sync_pegasus_qa_agents_lab.py → push to github.com/mgunjal11/pegasus-qa-agents-lab",
        ]
        for j, line in enumerate(lines):
            para = scripts.text_frame.paragraphs[0] if j == 0 else scripts.text_frame.add_paragraph()
            para.text = line
            para.font.size = Pt(11 if j else 12)
            para.font.bold = j == 0
            para.font.color.rgb = NAVY if j == 0 else BODY
            para.font.name = FONT

    def case_studies_slide(self):
        """Proven MSC outcomes — four pilot validations (featured: MSC-205625)."""
        s = self.blank()
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHITE
        self.footer(s)
        band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.28), W, Inches(1.05))
        band.fill.solid()
        band.fill.fore_color.rgb = NAVY
        band.line.fill.background()
        tit = s.shapes.add_textbox(Inches(0.55), Inches(0.42), Inches(8.5), Inches(0.45))
        tit.text_frame.text = "Proven MSC outcomes"
        tit.text_frame.paragraphs[0].font.size = Pt(26)
        tit.text_frame.paragraphs[0].font.bold = True
        tit.text_frame.paragraphs[0].font.color.rgb = WHITE
        tit.text_frame.paragraphs[0].font.name = FONT
        sub = s.shapes.add_textbox(Inches(0.55), Inches(0.88), Inches(9.5), Inches(0.35))
        sub.text_frame.text = f"Four production runs · featured story {PRIMARY_EXAMPLE_KEY} (this deck)"
        sub.text_frame.paragraphs[0].font.size = Pt(11)
        sub.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
        sub.text_frame.paragraphs[0].font.name = FONT
        feat = REPORT_MATRIX[PRIMARY_EXAMPLE_KEY]
        self._kpi_tile(s, feat["dev_code_pct"], "Featured dev code", Inches(9.5), Inches(0.42), Inches(1.35), Inches(0.85), SOFT_GOLD, NAVY)
        self._kpi_tile(s, feat["testplan_ac_pct"], "Featured plan AC", Inches(10.95), Inches(0.42), Inches(1.35), Inches(0.85), AMBER_BG, AMBER)
        cases = []
        for key in ("MSC-212571", "MSC-205625", "MSC-204417", "MSC-195138"):
            r = REPORT_MATRIX[key]
            notes = f"{r['pr_note']} · {r['testplan_note']}"
            if r.get("ladr_note") and r["ladr_note"] != "—":
                notes = f"{r['ladr_note']} · {notes}"
            cases.append(
                (
                    key,
                    r["type"],
                    r["dev_code_pct"],
                    r["dev_tests_pct"],
                    r["testplan_ac_pct"],
                    notes[:52] + ("…" if len(notes) > 52 else ""),
                )
            )
        headers = ["Ticket", "Type", "Code", "Dev tests", "Plan AC", "Notes"]
        col_w = Inches(2.03)
        for ci, h in enumerate(headers):
            left = Inches(0.55) + col_w * ci
            self._rect(s, left, Inches(1.48), col_w - Inches(0.04), Inches(0.4), NAVY, h, 9, True, WHITE)

        def _cell_color(ci: int, cell: str):
            if ci <= 1:
                return DARK
            if cell.endswith("%"):
                return GREEN if cell.startswith("100") else AMBER
            return BODY

        for ri, row in enumerate(cases):
            top = Inches(1.98) + Inches(0.55) * ri
            is_feat = row[0] == PRIMARY_EXAMPLE_KEY
            row_fill = SOFT_GOLD if is_feat else (LIGHT_GRAY if ri % 2 else WHITE)
            for ci, cell in enumerate(row):
                left = Inches(0.55) + col_w * ci
                bg, fg = self._cell_status_style(cell, default_bg=row_fill)
                label = f"★ {cell}" if is_feat and ci == 0 else cell
                self._rect(
                    s, left, top, col_w - Inches(0.04), Inches(0.5), bg,
                    label[:48] if ci == 5 else label[:20], 8, ci == 0,
                    fg if ci == 0 else _cell_color(ci, cell),
                )
            if is_feat:
                star = s.shapes.add_textbox(Inches(12.15), top + Inches(0.12), Inches(0.5), Inches(0.3))
                star.text_frame.text = "★"
                star.text_frame.paragraphs[0].font.size = Pt(14)
                star.text_frame.paragraphs[0].font.color.rgb = GOLD


    def closing_slide(self):
        s = self.blank()
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = NAVY
        bg.line.fill.background()
        c1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(4.5), Inches(4), Inches(4))
        c1.fill.solid()
        c1.fill.fore_color.rgb = NAVY_LIGHT
        c1.fill.transparency = 0.3
        c1.line.fill.background()
        t1 = s.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.3), Inches(0.8))
        t1.text_frame.text = "It's time to"
        t1.text_frame.paragraphs[0].font.size = Pt(36)
        t1.text_frame.paragraphs[0].font.color.rgb = WHITE
        t1.text_frame.paragraphs[0].font.name = FONT
        t1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        t2 = s.shapes.add_textbox(Inches(1.5), Inches(2.95), Inches(10.3), Inches(1.0))
        t2.text_frame.text = "Outcreate"
        t2.text_frame.paragraphs[0].font.size = Pt(52)
        t2.text_frame.paragraphs[0].font.bold = True
        t2.text_frame.paragraphs[0].font.color.rgb = GOLD
        t2.text_frame.paragraphs[0].font.name = FONT
        t2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        t3 = s.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9.3), Inches(0.5))
        t3.text_frame.text = f"/msc-dev-code-and-qa-test-coverage-validator {{YOUR-MSC-KEY}}  ·  Developed by {REPORT_DEVELOPER}"
        t3.text_frame.paragraphs[0].font.size = Pt(14)
        t3.text_frame.paragraphs[0].font.name = "Consolas"
        t3.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
        t3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.55), W, Inches(0.95))
        band.fill.solid()
        band.fill.fore_color.rgb = GOLD
        band.line.fill.background()

    def build_all(self):
        self.title_slide()
        self.agenda_slide()
        self.executive_summary()
        self.section_slide("01", "The challenge", "Why fragmented Jira · GitHub · Excel test plans · Confluence LADR blocks confident releases")
        self.challenge_slide()
        self.idea_generation_slide()
        self.section_slide("02", "The solution", "Cursor subagent · automated workflow · solution architecture")
        self.personas_slide()
        self.workflow_slide()
        self.architecture_slide()
        self.section_slide("03", "HTML report", "Eight sections · info icons · metrics · traceability")
        self.report_sections_slide()
        self.report_ui_slide()
        self.metrics_dashboard_slide()
        self.requirement_traceability_slide()
        self.confluence_ladr_slide()
        self.testplan_slide()
        self.dev_qa_slide()
        self.section_slide("04", "Enablement", "One-time setup · scripts · cache · permissions")
        self.setup_slide()
        self.section_slide("05", "Outcomes", "MSC case studies · latest report matrix · real coverage metrics")
        self.report_matrix_slide()
        self.case_studies_slide()
        self.closing_slide()


def build(
    out: Path,
    *,
    keep_prefix: int = KEEP_PREFIX_SLIDES,
    base_ppt: Path | None = None,
    report_html: Path | None = None,
) -> None:
    refresh_report_matrix_from_html()
    _set_latest_example(PRIMARY_EXAMPLE_KEY)
    report_path = report_html or ROOT / "reports" / REPORT_MATRIX[PRIMARY_EXAMPLE_KEY]["report_file"]
    if not report_path.is_file():
        report_path = DEFAULT_REPORT_HTML
    report_data = parse_msc_report_html(report_path)
    print(f"Report source: {report_path.name} ({report_data['issue_key']})")

    base = base_ppt or DEFAULT_BASE_PPT
    out.parent.mkdir(parents=True, exist_ok=True)
    if base.is_file() and keep_prefix > 0:
        prs = Presentation(str(base))
        while len(prs.slides) > keep_prefix:
            _delete_slide(prs, len(prs.slides) - 1)
        d = Deck(prs=prs, footer_start=keep_prefix)
        d.build_report_tail_from_html(report_data)
        prs.save(str(out))
        print(f"Kept slides 1–{keep_prefix} from {base.name}; rebuilt slides {keep_prefix + 1}–{len(prs.slides)} from HTML")
    else:
        print(f"Note: {base} not found — generating full deck from script")
        d = Deck()
        d.build_all()
        d.prs.save(str(out))
    print(out.resolve())


if __name__ == "__main__":
    import shutil
    import sys

    dest = Path(sys.argv[1]) if len(sys.argv) > 1 else ROOT / "docs" / "MSC-Dev-Code-and-QA-Test-Coverage-Validator-Guide.pptx"
    build(dest)
    reports_copy = ROOT / "reports" / "MSC-Dev-Code-and-QA-Test-Coverage-Validator-Guide.pptx"
    if dest.resolve() != reports_copy.resolve():
        try:
            shutil.copy2(dest, reports_copy)
            print(reports_copy.resolve())
        except OSError as exc:
            print(f"Note: could not copy to {reports_copy} ({exc})")
