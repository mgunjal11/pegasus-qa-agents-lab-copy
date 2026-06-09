"""PowerPoint slideshow hover tooltips via OOXML (python-pptx has no public screen_tip API)."""

from __future__ import annotations

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt

MAX_TOOLTIP_LEN = 255
VISIBLE_DESC_MAX = 110


def truncate_visible(text: str | None, max_chars: int = VISIBLE_DESC_MAX) -> str:
    if not text or not str(text).strip():
        return ""
    t = " ".join(str(text).split())
    if len(t) <= max_chars:
        return t
    return t[: max_chars - 1].rstrip() + "…"


def set_hover_tooltip(shape, text: str | None) -> None:
    """Show *text* when the pointer hovers the shape in Slide Show mode."""
    if not text or not str(text).strip():
        return
    tip = " ".join(str(text).split())[:MAX_TOOLTIP_LEN]
    el = shape._element
    c_nv = None
    if hasattr(el, "nvSpPr") and el.nvSpPr is not None:
        c_nv = el.nvSpPr.cNvPr
    elif hasattr(el, "nvPicPr") and el.nvPicPr is not None:
        c_nv = el.nvPicPr.cNvPr
    if c_nv is None:
        return
    for child in list(c_nv):
        if child.tag.endswith("}hlinkHover"):
            c_nv.remove(child)
    h = OxmlElement("a:hlinkHover")
    h.set("tooltip", tip)
    h.set("action", "ppaction://noaction")
    c_nv.append(h)


def add_visible_description(
    slide,
    left,
    top,
    width,
    height,
    text: str | None,
    *,
    font_pt: int = 8,
    font_rgb=None,
    font_name: str = "Calibri",
):
    """Always-visible caption (gray). Full *text* still available on hover."""
    from pptx.dml.color import RGBColor

    cap = truncate_visible(text)
    if not cap:
        return None
    font_rgb = font_rgb or RGBColor(0x64, 0x74, 0x8B)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = cap
    p = tf.paragraphs[0]
    p.font.size = Pt(font_pt)
    p.font.name = font_name
    p.font.color.rgb = font_rgb
    p.font.italic = True
    set_hover_tooltip(box, text)
    return box


def add_info_tip(
    slide,
    left,
    top,
    tooltip: str,
    *,
    size=None,
    fill_rgb=None,
    border_rgb=None,
    text_rgb=None,
    font_name: str = "Georgia",
    label: str = "i",
):
    """Visible *i* badge; hover shows full *tooltip* in Slide Show (F5)."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor

    if size is None:
        from pptx.util import Inches

        size = Inches(0.28)
    fill_rgb = fill_rgb or RGBColor(0xFF, 0x5E, 0x4F)
    border_rgb = border_rgb or RGBColor(0xE2, 0x8F, 0x87)
    text_rgb = text_rgb or RGBColor(0xFF, 0xFF, 0xFF)

    o = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    o.fill.solid()
    o.fill.fore_color.rgb = fill_rgb
    o.line.color.rgb = border_rgb
    o.line.width = Pt(0.5)
    o.text_frame.text = label
    p = o.text_frame.paragraphs[0]
    p.font.size = Pt(9)
    p.font.italic = True
    p.font.bold = True
    p.font.name = font_name
    p.font.color.rgb = text_rgb
    p.alignment = PP_ALIGN.CENTER
    o.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_hover_tooltip(o, tooltip)
    return o
